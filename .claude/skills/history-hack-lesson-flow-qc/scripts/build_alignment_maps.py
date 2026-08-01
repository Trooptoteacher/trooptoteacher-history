"""
build_alignment_maps.py — extract the deck ↔ workbook alignment maps and run the mechanical
cross-checks that the lesson-flow QC agent verifies visually.

This does the *programmatic* half of the student-journey audit: it builds three maps and flags
the objective mismatches (missing slide refs, DI-count drift, role-order drift, orphan cues). The
agent still renders and READS the slides/pages to confirm content parity and pedagogical flow —
this script only finds the leads. Text is authoritative for structure; pixels are authoritative
for content (see SKILL.md).

Inputs (positional): TEACHER_DECK.pptx  STUDENT_DECK.pptx  STUDENT_WORKBOOK.docx
Output: prints a per-standard report + a FLAGS list (BLOCKER/MAJOR/MINOR) to stdout as JSON-ish text.

Grounded in the Course Standard file conventions:
- Teacher deck slide titles: "US.xx · <ROLE>" (e.g. "US.45 · DIRECT INSTRUCTION").
- Student deck: bare "US.xx" starts a standard block; role slides are titled by keyword
  (DIRECT INSTRUCTION / SOURCE IT FIRST / THREE PERSPECTIVES / KEY VOCABULARY / PROGRESS CHECK);
  on-slide cue "✍ In your workbook · <activity>".
- Workbook activity headers: "Activity N — <name> — US.xx ... ▶ Deck slide(s) N".
- Guided Cornell cue column: "▶ Deck · DI N of M" (M = DI segments the workbook expects).
"""

import re
import sys
from collections import defaultdict

from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation

STD = re.compile(r"US\.(\d{2})")
DECK_REF = re.compile(r"▶\s*Deck slides?\s*([\d]+)(?:\s*[–-]\s*(\d+))?")
DI_OF = re.compile(r"DI\s*(\d+)\s*of\s*(\d+)")
ROLE_KEYS = ["DIRECT INSTRUCTION", "SOURCE IT FIRST", "PRIMARY SOURCE", "THREE PERSPECTIVES",
             "KEY VOCABULARY", "PROGRESS CHECK", "CHECK FOR UNDERSTANDING", "HOOK",
             "QUICK REVIEW", "CONFIDENCE CHECK", "PEOPLE WHO SHAPED", "STUDENT ACTIVITY"]


def _slide_title(s):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().split("\n")[0]
    return ""


def _slide_text(s):
    parts = [sh.text_frame.text for sh in s.shapes if sh.has_text_frame]
    if s.has_notes_slide:
        parts.append("NOTES:" + s.notes_slide.notes_text_frame.text)
    return "\n".join(parts)


def _role(title):
    up = title.upper()
    for k in ROLE_KEYS:
        if k in up:
            return k
    return "OTHER"


def map_teacher_deck(path):
    """standard -> ordered [(slide_no, role, has_write_cue)]"""
    p = Presentation(path)
    m = defaultdict(list)
    for i, s in enumerate(p.slides, 1):
        t = _slide_title(s)
        mm = STD.search(t)
        if not mm:
            continue
        code = "US." + mm.group(1)
        txt = _slide_text(s)
        m[code].append((i, _role(t), "✍" in txt))
    return m


def map_student_deck(path):
    """standard -> ordered [(slide_no, role, cue_activity_or_None)]. Blocks start at a bare 'US.xx' slide."""
    p = Presentation(path)
    m = defaultdict(list)
    cur = None
    for i, s in enumerate(p.slides, 1):
        t = _slide_title(s).strip()
        txt = _slide_text(s)
        mm = STD.fullmatch(t) or (STD.fullmatch(t.split()[0]) if t.split() else None)
        if STD.fullmatch(t):            # a bare "US.xx" title slide starts a block
            cur = "US." + STD.search(t).group(1)
            m[cur].append((i, "TITLE", None))
            continue
        if cur is None:
            continue
        cue = None
        cm = re.search(r"✍ In your workbook · ([^\n]+)", txt)
        if cm:
            cue = cm.group(1).strip()
        m[cur].append((i, _role(t), cue))
    return m


def map_workbook(path):
    """standard -> {activities:[(n,name,ref_lo,ref_hi)], di_expected:int}"""
    doc = Document(path)
    body = doc.element.body

    def ptext(el):
        return "".join(n.text or "" for n in el.iter(qn("w:t")))

    m = defaultdict(lambda: {"activities": [], "di_expected": 0})
    for ch in body.iter(qn("w:p")):
        t = ptext(ch)
        am = re.match(r"Activity\s+(\d+)\s+[—-]\s+(.+?)\s+[—-]\s+US\.(\d{2})", t)
        if am:
            code = "US." + am.group(3)
            ref = DECK_REF.search(t)
            lo = int(ref.group(1)) if ref else None
            hi = int(ref.group(2)) if (ref and ref.group(2)) else lo
            m[code]["activities"].append((int(am.group(1)), am.group(2).strip(), lo, hi))
        di = DI_OF.search(t)
        if di:
            code_m = STD.search(t)
            # DI cue lives in the Cornell table cell; attribute to the nearest standard seen on the line
            if code_m:
                m["US." + code_m.group(1)]["di_expected"] = max(
                    m["US." + code_m.group(1)]["di_expected"], int(di.group(2)))
    return m


def audit(teacher, student, workbook):
    T = map_teacher_deck(teacher)
    S = map_student_deck(student)
    W = map_workbook(workbook)
    student_slidecount = len(Presentation(student).slides._sldIdLst)
    flags = []
    codes = sorted(set(T) | set(S) | set(W), key=lambda c: int(c.split(".")[1]))
    for c in codes:
        t_di = sum(1 for _, r, _ in T.get(c, []) if r == "DIRECT INSTRUCTION")
        s_di = sum(1 for _, r, _ in S.get(c, []) if r == "DIRECT INSTRUCTION")
        w_di = W.get(c, {}).get("di_expected", 0)
        print(f"\n{c}: teacher DI={t_di}  student DI={s_di}  workbook expects DI={w_di}")
        print(f"  teacher: {[ (n,r) for n,r,_ in T.get(c,[]) ]}")
        print(f"  student: {[ (n,r) for n,r,_ in S.get(c,[]) ]}")
        print(f"  workbook activities: {[ (n,nm,lo,hi) for n,nm,lo,hi in W.get(c,{}).get('activities',[]) ]}")

        # CHECK 1 — DI coverage: student deck must cover every DI segment the workbook keys to
        if w_di and s_di < w_di:
            flags.append(("BLOCKER", c, f"Guided Cornell keys DI 1..{w_di} but student deck has only "
                          f"{s_di} DIRECT INSTRUCTION slide(s) — 'DI {s_di+1} of {w_di}' has no home in the student's own deck."))
        if t_di and s_di and s_di != t_di:
            flags.append(("MAJOR", c, f"Teacher deck teaches {t_di} DI segments but the student deck has {s_di} — "
                          f"student review deck does not 100% cover what was taught."))
        # CHECK 2 — every workbook ▶ Deck ref resolves within the student deck
        for n, nm, lo, hi in W.get(c, {}).get("activities", []):
            if lo and (lo > student_slidecount or (hi and hi > student_slidecount)):
                flags.append(("BLOCKER", c, f"Activity {n} ({nm}) points to ▶ Deck slide {lo}"
                              f"{'–'+str(hi) if hi and hi!=lo else ''} but the student deck has only {student_slidecount} slides."))
        # CHECK 3 — vocab-before-content order (student deck): vocab should not trail all content
        roles = [r for _, r, _ in S.get(c, [])]
        if "KEY VOCABULARY" in roles and "DIRECT INSTRUCTION" in roles:
            if roles.index("KEY VOCABULARY") > roles.index("DIRECT INSTRUCTION"):
                flags.append(("MAJOR", c, "Student deck presents KEY VOCABULARY *after* DIRECT INSTRUCTION, "
                              "but the workbook does Vocabulary as Activity 1–2 (before Cornell notes) — sequence mismatch."))
        # CHECK 4 — orphan write cues (student deck cue with no matching workbook activity name)
        wb_names = " ".join(nm.lower() for _, nm, _, _ in W.get(c, {}).get("activities", []))
        for n, r, cue in S.get(c, []):
            if cue:
                key = cue.split("(")[0].strip().lower().split("&")[0].strip()
                if key and wb_names and key.split()[0] not in wb_names:
                    flags.append(("MINOR", c, f"Student deck slide {n} cue '✍ {cue}' has no clearly matching workbook activity."))

    print("\n\n===== FLAGS (most severe first) =====")
    order = {"BLOCKER": 0, "MAJOR": 1, "MINOR": 2}
    for sev, c, msg in sorted(flags, key=lambda f: (order[f[0]], f[1])):
        print(f"[{sev}] {c}: {msg}")
    print(f"\nTOTAL: {sum(1 for f in flags if f[0]=='BLOCKER')} blocker, "
          f"{sum(1 for f in flags if f[0]=='MAJOR')} major, {sum(1 for f in flags if f[0]=='MINOR')} minor")
    return flags


if __name__ == "__main__":
    if len(sys.argv) != 4:
        print(__doc__)
        sys.exit(1)
    audit(sys.argv[1], sys.argv[2], sys.argv[3])
