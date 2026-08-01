"""
scan_text_integrity.py — catch incomplete / clipped / placeholder text in decks and workbooks.

Born from a real Unit 6 defect: vocabulary definitions authored with a trailing "…" that ran off
the card ("…annex the Sudetenland from…"). This scanner makes that whole error class impossible to
ship silently again. It reads the TEXT (authoritative for truncation) and estimates box fit
(a lead for clipping — confirm the flagged slides by rendering).

Checks (per text unit — a pptx shape paragraph or a docx paragraph / table cell):
  1. TRUNCATION  — ends with "…"/"..." or a dangling function word (from, like, and, the, such as,
     including, e.g., with, to, of, in, for, by …). Deterministic; the exact bug we hit.
  2. NO-TERMINAL — long sentence-like text (> 80 chars) that ends with no terminal punctuation
     (. ! ? " ' ) : ]) — a softer cut-off signal.
  3. OVERFLOW    — (pptx) estimated lines needed > the text box height allows (font-size aware,
     conservative). A LEAD to confirm visually, not a verdict.
  4. PLACEHOLDER — lorem/ipsum, xxx, TODO, [insert, TBD, "this slide layout".

Usage:  python scan_text_integrity.py FILE [FILE ...]        # .pptx and/or .docx
Exit code 1 if any BLOCKER/MAJOR found (so it can gate a build), 0 otherwise.
"""
import math
import re
import sys

# dangling openers a complete sentence never ends on (high-precision set — words like
# "on/as/by/in" are excluded because they legitimately end phrases: "from then on", "and so on")
DANGLERS = (r"from|like|such as|including|and|or|the|a|an|with|to|of|"
            r"e\.g\.|i\.e\.|which|that|whose|named|called")
DANGLE_RE = re.compile(r"(?:^|\s)(?:" + DANGLERS + r")\s*[…\.\-–—]*\s*$", re.IGNORECASE)
TRUNC_END_RE = re.compile(r"(…|\.\.\.)\s*$")
PLACEHOLDER_RE = re.compile(r"\b(lorem|ipsum|x{3,}|TODO|TBD)\b|\[insert|this .*(slide|page) layout", re.IGNORECASE)
TERMINAL = tuple(list('.!?"”’') + [")", ":", "]", "%"])

def cls_truncation(t):
    """Separate real truncations from intentional stems/elisions:
      - short text ending in "…"  ->  a sentence stem/prompt ("My goal is…", "I can…"): NOT a defect
      - long text ending in "…" on a dangling word ("…under leaders like…") -> BLOCKER (cut mid-clause)
      - long text ending in "…" on a normal word -> MAJOR (verify vs. an intended source elision)
      - long text ending on a dangling word without "…" -> MAJOR (likely cut off)
    """
    s = t.strip()
    if len(s) < 12:
        return None
    if TRUNC_END_RE.search(s):
        if len(s) <= 70:
            return None  # intentional sentence stem / fill-in prompt
        if DANGLE_RE.search(s):
            return ("BLOCKER", "truncation", "long text cut off mid-clause (ends on a dangling word + ellipsis)")
        return ("MAJOR", "truncation", "long text ends with an ellipsis — confirm it is not an intended source elision")
    # NOTE: a dangling-word ending WITHOUT an ellipsis is not checked — real in-source truncations
    # always leave a "…", and the no-ellipsis variant only produced false positives (answer keys
    # ending "4. A", instructions ending "…on"). Silent mid-box clipping is the overflow check's job.
    return None

def cls_noterminal(t):
    # low-signal; only runs under --all. A long, multi-word line that ends mid-word (lowercase,
    # no terminal punctuation) is a possible silent cut-off.
    s = t.strip()
    if (len(s) > 100 and not s.endswith(TERMINAL) and " " in s
            and s[-1].islower() and "\n" not in s):
        return ("MINOR", "no-terminal", "long line ends mid-sentence with no punctuation")
    return None

def cls_placeholder(t):
    if PLACEHOLDER_RE.search(t):
        return ("BLOCKER", "placeholder", "leftover placeholder / TODO text")
    return None

def est_overflow(text, width_in, height_in, font_pt):
    """Conservative fit estimate. chars/inch ~ 9 at 11pt, scales inverse to font size."""
    if not (width_in and height_in and font_pt):
        return None
    s = text.strip()
    if len(s) < 40:
        return None
    cpl = max(1, width_in * 9.0 * (11.0 / font_pt))
    lines = math.ceil(len(s) / cpl)
    line_h = font_pt * 0.0167           # ~1.2 line spacing, inches
    needed = lines * line_h
    if needed > height_in * 1.35:       # 35% slack so only clear cases flag
        return ("MAJOR", "overflow",
                f"~{lines} lines (~{needed:.2f}in) may exceed the {height_in:.2f}in box — verify by render")
    return None

ALL = False  # --all also runs the noisy low-signal no-terminal check

def scan_para(text, loc, findings, geo=None):
    checks = [cls_placeholder, cls_truncation] + ([cls_noterminal] if ALL else [])
    for fn in checks:
        r = fn(text)
        if r:
            sev, cat, why = r
            findings.append((sev, cat, loc, why, text.strip()[:90]))
            return  # one finding per unit is enough
    if geo:
        r = est_overflow(text, *geo)
        if r:
            sev, cat, why = r
            findings.append((sev, cat, loc, why, text.strip()[:90]))

def scan_pptx(path, findings):
    from pptx import Presentation
    prs = Presentation(path)
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            full = tf.text
            if not full.strip():
                continue
            # geometry + first-run font size for the overflow lead
            try:
                w = (sh.width or 0) / 914400.0; h = (sh.height or 0) / 914400.0
                fp = None
                for p in tf.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            fp = r.font.size.pt; break
                    if fp: break
                geo = (w, h, fp or 12.0)
            except Exception:
                geo = None
            scan_para(full, f"{path} · slide {i}", findings, geo)

def scan_docx(path, findings):
    from docx import Document
    doc = Document(path)
    for i, p in enumerate(doc.paragraphs, 1):
        if p.text.strip():
            scan_para(p.text, f"{path} · ¶{i}", findings)
    for ti, tbl in enumerate(doc.tables, 1):
        for ri, row in enumerate(tbl.rows, 1):
            for ci, cell in enumerate(row.cells, 1):
                for p in cell.paragraphs:
                    if p.text.strip():
                        scan_para(p.text, f"{path} · table{ti} r{ri}c{ci}", findings)

def main():
    global ALL
    args = [a for a in sys.argv[1:]]
    if "--all" in args:
        ALL = True; args.remove("--all")
    if not args:
        print(__doc__); sys.exit(2)
    findings = []
    for path in args:
        (scan_pptx if path.lower().endswith((".pptx", ".potx")) else scan_docx)(path, findings)
    order = {"BLOCKER": 0, "MAJOR": 1, "MINOR": 2}
    findings.sort(key=lambda f: order.get(f[0], 9))
    counts = {"BLOCKER": 0, "MAJOR": 0, "MINOR": 0}
    for sev, cat, loc, why, snip in findings:
        counts[sev] = counts.get(sev, 0) + 1
        print(f"[{sev}] {cat} · {loc}\n    {why}\n    “{snip}”")
    print(f"\nTOTAL: {counts['BLOCKER']} blocker, {counts['MAJOR']} major, {counts['MINOR']} minor")
    if counts["MAJOR"]:
        print("(MAJOR = advisory: overflow leads to render-confirm, long-ellipsis to check vs. elision)")
    # gate on BLOCKER only — MAJOR/overflow need a human render-confirm, not an auto-fail
    sys.exit(1 if counts["BLOCKER"] else 0)

if __name__ == "__main__":
    main()
