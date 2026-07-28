#!/usr/bin/env python3
"""Insert the alignment layer into the AUTHENTIC teacher deck (any unit).

Front: 1 intro ("How This Deck Drives the Workbook").
Per standard: 1 "LESSON -> WORKBOOK MAP" slide placed right AFTER that standard's
STANDARD divider, so it opens each standard's cycle. Authentic slides and every
already-embedded image pass through untouched (only the layer slides are cloned in,
with their — image-free — shapes).

Usage:
  python3 merge_align.py TEACHER_DECK.pptx OUT.pptx align.json teacher_align_layer.pptx
"""
import copy, json, sys
from io import BytesIO
from pptx import Presentation
from pptx.oxml.ns import qn

TGT, OUT, ALIGN_JSON, LAYER = sys.argv[1], sys.argv[2], sys.argv[3], sys.argv[4]
ALIGN = json.load(open(ALIGN_JSON))

tgt = Presentation(TGT)
layer = Presentation(LAYER)
blank = tgt.slide_layouts[6] if len(tgt.slide_layouts) > 6 else tgt.slide_layouts[-1]
SHAPE_TAGS = {qn('p:sp'), qn('p:pic'), qn('p:graphicFrame'), qn('p:grpSp'), qn('p:cxnSp')}
sldIdLst = tgt.slides._sldIdLst
orig = list(sldIdLst)
NSRC = len(orig)

def add_from(src_prs, idx):
    src_slide = src_prs.slides[idx]
    slide = tgt.slides.add_slide(blank)
    spTree = slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for sp in list(spTree):
        if sp.tag in SHAPE_TAGS:
            spTree.remove(sp)
    src_cSld = src_slide._element.find(qn('p:cSld'))
    bg = src_cSld.find(qn('p:bg'))
    if bg is not None:
        slide._element.find(qn('p:cSld')).insert(0, copy.deepcopy(bg))
    for sp in list(src_cSld.find(qn('p:spTree'))):
        if sp.tag not in SHAPE_TAGS:
            continue
        new_sp = copy.deepcopy(sp)
        for blip in new_sp.findall('.//' + qn('a:blip')):
            rId = blip.get(qn('r:embed'))
            if not rId:
                continue
            blob = src_slide.part.related_part(rId).blob
            _, new_rId = slide.part.get_or_add_image_part(BytesIO(blob))
            blip.set(qn('r:embed'), new_rId)
        spTree.append(new_sp)
    return sldIdLst[-1]

intro = add_from(layer, 0)
maps = {rec["code"]: add_from(layer, 1 + i) for i, rec in enumerate(ALIGN["recs"])}
div_after = {rec["divider_slide"]: rec["code"] for rec in ALIGN["recs"]}

order = [intro]
for i, el in enumerate(orig, start=1):        # i = 1-based ORIGINAL slide number
    order.append(el)
    if i in div_after:                        # drop the map in right after the divider
        order.append(maps[div_after[i]])

assert len(order) == NSRC + 1 + len(ALIGN["recs"]), (len(order), NSRC)
for el in list(sldIdLst):
    sldIdLst.remove(el)
for el in order:
    sldIdLst.append(el)

tgt.save(OUT)
print("SAVED", OUT, "slides:", len(tgt.slides._sldIdLst),
      "(src", NSRC, "+", 1 + len(ALIGN["recs"]), "inserted)")
