#!/usr/bin/env python3
# Teacher-deck variant of deck_merge: blocks start at the FIRST slide of each
# US.NN run (e.g. "US.45 · QUICK REVIEW"), and any pre-first-block front matter
# is preserved (intro platinum slide inserted right after it).
import sys, copy, re
from io import BytesIO
from pptx import Presentation
from pptx.oxml.ns import qn
src_deck, udl_deck, out = sys.argv[1], sys.argv[2], sys.argv[3]
tgt=Presentation(src_deck); udl=Presentation(udl_deck)
blank=tgt.slide_layouts[6] if len(tgt.slide_layouts)>6 else tgt.slide_layouts[-1]
SHAPE_TAGS={qn('p:sp'),qn('p:pic'),qn('p:graphicFrame'),qn('p:grpSp'),qn('p:cxnSp')}
def ft(s):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip(): return sh.text_frame.text.strip().splitlines()[0].strip()
    return ""
def add_from(idx):
    s=udl.slides[idx]; sl=tgt.slides.add_slide(blank)
    tcs=sl._element.find(qn('p:cSld')); tree=tcs.find(qn('p:spTree'))
    for sp in list(tree):
        if sp.tag in SHAPE_TAGS: tree.remove(sp)
    scs=s._element.find(qn('p:cSld')); bg=scs.find(qn('p:bg'))
    if bg is not None: tcs.insert(0,copy.deepcopy(bg))
    for sp in list(scs.find(qn('p:spTree'))):
        if sp.tag not in SHAPE_TAGS: continue
        n=copy.deepcopy(sp)
        for blip in n.findall('.//'+qn('a:blip')):
            rId=blip.get(qn('r:embed'))
            if rId:
                blob=s.part.related_part(rId).blob
                _,nr=sl.part.get_or_add_image_part(BytesIO(blob)); blip.set(qn('r:embed'),nr)
        tree.append(n)
    return list(tgt.slides._sldIdLst)[-1]
sldIdLst=tgt.slides._sldIdLst
orig=list(sldIdLst)
# detect block starts (first slide of each code run)
blocks=[]; prev=None
for i,s in enumerate(tgt.slides):
    m=re.match(r'(US\.\d+)',ft(s)); code=m.group(1) if m else None
    if code and code!=prev: blocks.append((i,code)); prev=code
assert blocks, "no standard blocks found"
codes=[c for _,c in blocks]; assert codes==sorted(codes),codes
first_start=blocks[0][0]
# clone layer slides
intro=add_from(0)
voc={c:add_from(1+2*i) for i,c in enumerate(codes)}
prog={c:add_from(2+2*i) for i,c in enumerate(codes)}
# order: front matter + intro + per block(divider..rest, vocab, progress)
order=list(orig[:first_start])+[intro]
for bi,(start,code) in enumerate(blocks):
    end=blocks[bi+1][0] if bi+1<len(blocks) else first_start+ (len(orig)-first_start)
    end=blocks[bi+1][0] if bi+1<len(blocks) else len(orig)
    order.extend(orig[start:end]); order.append(voc[code]); order.append(prog[code])
expected=len(orig)+1+2*len(codes)
assert len(order)==expected,(len(order),expected)
for el in list(sldIdLst): sldIdLst.remove(el)
for el in order: sldIdLst.append(el)
tgt.save(out)
print("SAVED",out,"slides:",len(tgt.slides._sldIdLst),"| standards:",len(codes),"| front-matter kept:",first_start)
