import sys, copy, re
from io import BytesIO
from pptx import Presentation
from pptx.oxml.ns import qn
SRC="Unit10_Teacher_Deck_ORIGINAL.pptx"
ENR="/tmp/claude-0/-home-user-trooptoteacher-history/c27aff16-3c79-5004-bf4a-470896c27424/scratchpad/Unit10_WordWall_Enrichment.pptx"
OUT="Unit10_Teacher_Deck_WW.pptx"
tgt=Presentation(SRC); enr=Presentation(ENR)
blank=tgt.slide_layouts[6] if len(tgt.slide_layouts)>6 else tgt.slide_layouts[-1]
SHAPE_TAGS={qn('p:sp'),qn('p:pic'),qn('p:graphicFrame'),qn('p:grpSp'),qn('p:cxnSp')}
def first_text(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().splitlines()[0].strip()
    return ""
def add_from(src_prs, idx):
    src=src_prs.slides[idx]; slide=tgt.slides.add_slide(blank)
    tcs=slide._element.find(qn('p:cSld')); tree=tcs.find(qn('p:spTree'))
    for sp in list(tree):
        if sp.tag in SHAPE_TAGS: tree.remove(sp)
    scs=src._element.find(qn('p:cSld')); bg=scs.find(qn('p:bg'))
    if bg is not None: tcs.insert(0,copy.deepcopy(bg))
    for sp in list(scs.find(qn('p:spTree'))):
        if sp.tag not in SHAPE_TAGS: continue
        n=copy.deepcopy(sp)
        for blip in n.findall('.//'+qn('a:blip')):
            rId=blip.get(qn('r:embed'))
            if rId:
                blob=src.part.related_part(rId).blob
                _,nr=slide.part.get_or_add_image_part(BytesIO(blob)); blip.set(qn('r:embed'),nr)
        tree.append(n)
    return slide
# map enrichment slides -> code
enr_by_code={}
for i,s in enumerate(enr.slides):
    m=re.match(r'(US\.\d+)\s*[·•]\s*KEY VOCAB',first_text(s))
    if m: enr_by_code[m.group(1)]=i
assert len(enr_by_code)==13, enr_by_code
# clone all 13 enriched slides into tgt (appended at tail); record sldId element
new_sld={}
sldIdLst=tgt.slides._sldIdLst
for code,idx in enr_by_code.items():
    add_from(enr, idx)
    new_sld[code]=list(sldIdLst)[-1]
# classify originals
orig=list(sldIdLst)[:210]   # first 210 are the source slides
def slide_of(el):
    rId=el.get(qn('r:id')); return tgt.slides._sldIdLst  # not used
# Build mapping index->(code, is_kv, is_us84_guided) using slide order
info=[]
for i,s in enumerate(tgt.slides):
    if i>=210: break
    ft=first_text(s)
    m=re.match(r'(US\.\d+)',ft)
    code=m.group(1) if m else None
    info.append((code,'KEY VOCAB' in ft.upper(),(code=='US.84' and 'GUIDED PRACTICE' in ft.upper())))
# new order
order=[]
for i,el in enumerate(orig):
    code,is_kv,is_g84=info[i]
    if is_kv:
        order.append(new_sld[code])       # replace old WW with enriched
    else:
        if is_g84:
            order.append(new_sld['US.84']) # insert US.84 WW before its guided practice
        order.append(el)
# sanity: all 13 enriched placed, all 12 old WW dropped
placed=[e for e in order if e in set(new_sld.values())]
assert len(placed)==13, len(placed)
# rewrite sldIdLst to the new order (drops old WWs + tail duplicates)
for el in list(sldIdLst): sldIdLst.remove(el)
for el in order: sldIdLst.append(el)
tgt.save(OUT)
print("SAVED",OUT,"slides:",len(tgt.slides._sldIdLst),"(was 210, +1 US.84 WW =",210+1,")")
