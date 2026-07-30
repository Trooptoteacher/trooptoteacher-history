# -*- coding: utf-8 -*-
"""Unit 1 Student (Lean) Deck — PLATINUM edit-in-place.
Reads the WORKING copy (real images intact), ADDS legend + per-standard
KEY VOCABULARY/LANGUAGE SUPPORT + PROGRESS CHECK (answer hidden, de-biased) +
Representation close, reorders, renumbers to file position, applies 9pt font
floor, descriptive alt text, and scans for fixed-track labels. Originals untouched."""
import copy, re, zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import unit1_data as D

IN  = 'working/Unit1_Student_Deck_WORKING.pptx'
OUT = 'deliverables/Unit1_Student_Deck_CourseStandard.pptx'

# palette
NAVY='1A2332'; GOLD='C9A84C'; LIGHT='C7D4E0'; WHITE='FFFFFF'
INK='1F2430'; GREEN='1B5E20'; RED='C62828'; BLUE='002858'; GREY='5C6470'
FOOT = D.UNIT['footer']
def C(h): return RGBColor.from_string(h)

prs = Presentation(IN)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[0]

def set_bg(slide, hexc=NAVY):
    cSld = slide._element.find(qn('p:cSld'))
    for ex in cSld.findall(qn('p:bg')): cSld.remove(ex)
    bg = cSld.makeelement(qn('p:bg'), {}); pr = bg.makeelement(qn('p:bgPr'), {})
    fill = pr.makeelement(qn('a:solidFill'), {}); clr = fill.makeelement(qn('a:srgbClr'), {'val':hexc})
    fill.append(clr); pr.append(fill); pr.append(pr.makeelement(qn('a:effectLst'), {}))
    bg.append(pr); cSld.insert(0, bg)

def new_slide():
    s = prs.slides.add_slide(BLANK)
    for sh in list(s.shapes): sh._element.getparent().remove(sh._element)
    set_bg(s)
    return s

def rect(s, l, t, w, h, fill=None, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = C(line); sp.line.width = Pt(lw)
    return sp

def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, sp_after=2):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ('left','right','top','bottom'): setattr(tf, 'margin_'+m, Emu(0))
    if isinstance(runs[0], tuple): runs = [runs]
    for pi, para in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        for (t_, fn, sz, col, bold) in para:
            r = p.add_run(); r.text = t_; r.font.name = fn; r.font.size = Pt(sz)
            r.font.color.rgb = C(col); r.font.bold = bold
    return tb

def chip(s, l, t, w, label, fill, fg=WHITE, sz=11):
    rect(s, l, t, w, 0.4, fill=fill)
    txt(s, l, t, w, 0.4, [(label, 'Trebuchet MS', sz, fg, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def footer(s, page):
    txt(s, 0.5, 7.08, 11.0, 0.3, [(FOOT, 'Trebuchet MS', 9, LIGHT, False)])   # 9pt floor
    txt(s, 12.23, 7.08, 0.6, 0.3, [(str(page), 'Trebuchet MS', 9, LIGHT, False)], align=PP_ALIGN.RIGHT)

def std_tag(s, code, l=11.6, t=0.35):
    rect(s, l, t, 1.15, 0.4, fill=GOLD)
    txt(s, l, t, 1.15, 0.4, [(code, 'Trebuchet MS', 12, NAVY, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------- LEGEND ----------
def build_legend():
    s = new_slide()
    txt(s, 0.5, 0.45, 12.3, 0.7, [('How This Deck Works — UDL & MTSS', 'Georgia', 30, WHITE, True)])
    txt(s, 0.5, 1.18, 12.3, 0.5, [('Universal supports built into every standard. One firm learning goal; flexible ways to reach it.', 'Georgia', 15, LIGHT, False)])
    items = [
        ('CORE PATH', GOLD, NAVY, 'The essential instruction every student receives.'),
        ('SUPPORT OPTION', GREEN, WHITE, 'An optional scaffold that keeps the goal, not lowers it.'),
        ('LANGUAGE SUPPORT', BLUE, WHITE, 'Vocabulary, pronunciation, and Spanish cognates for access.'),
        ('RESPONSE CHOICE', RED, WHITE, 'Show learning by writing, saying/recording, or diagramming.'),
        ('PROGRESS CHECK', '6A4E9C', WHITE, 'A quick DOK-2/3 check to guide reteach or extend.'),
        ('EXTENSION', '2F6F4E', WHITE, 'A deeper challenge once the goal is met.'),
    ]
    x0, y0, w, h, gx, gy = 0.5, 1.95, 3.95, 1.15, 0.22, 0.28
    for i, (lab, fill, fg, gloss) in enumerate(items):
        col, row = i % 3, i // 3
        l = x0 + col * (w + gx); t = y0 + row * (h + gy)
        rect(s, l, t, w, h, fill=WHITE)
        chip(s, l, t, w, lab, fill, fg=fg, sz=12)
        txt(s, l + 0.14, t + 0.5, w - 0.28, h - 0.55, [(gloss, 'Calibri', 12, INK, False)])
    rect(s, 0.5, 5.95, 12.3, 0.95, fill='232F40', line=GOLD, lw=1.5)
    txt(s, 0.7, 6.08, 11.9, 0.8, [
        [('Options are available by design and support access without lowering the learning goal.  ', 'Calibri', 12.5, WHITE, True)],
        [('Supports work alongside — never in place of — required IEP or 504 accommodations.', 'Calibri', 12.5, LIGHT, False)],
    ])
    footer(s, 0)
    return s

# ---------- VOCAB ----------
def build_vocab(code):
    s = new_slide()
    chip(s, 0.5, 0.3, 4.7, 'KEY VOCABULARY · LANGUAGE SUPPORT', BLUE, sz=11)
    std_tag(s, code)
    txt(s, 0.5, 0.9, 12.3, 0.5, [('Say it, see it, connect it — terms you need for this standard.', 'Georgia', 15, LIGHT, False)])
    terms = D.VOCAB[code]
    n = len(terms); w = (12.3 - 0.4 * (n - 1)) / n
    for i, (term, say, es, deff) in enumerate(terms):
        l = 0.5 + i * (w + 0.4)
        rect(s, l, 1.6, w, 4.9, fill=WHITE)
        rect(s, l, 1.6, w, 0.75, fill=NAVY)
        txt(s, l + 0.14, 1.68, w - 0.28, 0.62, [(term, 'Georgia', 15, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
        y = 2.5
        if say:
            txt(s, l + 0.14, y, w - 0.28, 0.3, [('say: ' + say, 'Trebuchet MS', 11, GREY, False)]); y += 0.36
        txt(s, l + 0.14, y, w - 0.28, 0.3, [('ES: ' + es, 'Trebuchet MS', 11, '8A6D1A', True)]); y += 0.42
        txt(s, l + 0.14, y, w - 0.28, 4.9 - (y - 1.6) - 0.2, [(deff, 'Calibri', 12, INK, False)])
    footer(s, 0)
    return s

# ---------- PROGRESS CHECK (answer hidden) ----------
def build_progress(code):
    item = D.CFU[code]; opts, _ = D.debiased_options(code)
    s = new_slide()
    chip(s, 0.5, 0.3, 3.4, 'PROGRESS CHECK', '6A4E9C', sz=11)
    chip(s, 4.0, 0.3, 1.5, 'DOK ' + str(item['dok']), GOLD, fg=NAVY, sz=11)
    std_tag(s, code)
    txt(s, 0.5, 0.95, 12.3, 1.15, [(item['stem'], 'Georgia', 18, WHITE, True)])
    y = 2.35
    for L in ['A', 'B', 'C', 'D']:
        rect(s, 0.5, y, 12.3, 0.92, fill='F2EFE6')
        rect(s, 0.5, y, 0.7, 0.92, fill=NAVY)
        txt(s, 0.5, y, 0.7, 0.92, [(L, 'Georgia', 18, GOLD, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, 1.4, y, 11.2, 0.92, [(opts[L], 'Calibri', 13.5, INK, False)], anchor=MSO_ANCHOR.MIDDLE)
        y += 1.02
    txt(s, 0.5, 6.5, 12.3, 0.4, [('RESPONSE CHOICE — commit before advancing: write it, say/record it, or diagram it.  Teacher reveal follows.', 'Trebuchet MS', 11, LIGHT, True)])
    footer(s, 0)
    return s

# ---------- REPRESENTATION ----------
def build_rep():
    s = new_slide()
    chip(s, 0.5, 0.3, 3.9, 'REACHED — AND LEFT OUT', GOLD, fg=NAVY, sz=11)
    txt(s, 0.5, 0.9, 12.3, 0.65, [('Whose story the industrial era told — and who paid its costs', 'Georgia', 22, WHITE, True)])
    txt(s, 0.5, 1.55, 12.3, 0.5, [('History Hack-authored synthesis grounded in this unit’s sourced record — read across the standards, not around them.', 'Georgia', 13, LIGHT, False)])
    cards = [
        ('American Indian Nations', GREEN, 'Westward expansion and the railroad (US.01) crossed treaty lands; reservation, assimilation, boarding-school, and Dawes Act policy (US.02) cut tribal holdings from 138 to 48 million acres by 1934.'),
        ('Chinese Laborers', BLUE, 'Chinese workers built the most dangerous stretches of the Central Pacific (US.01) yet faced wage discrimination and, in 1882, the first federal law to bar a nationality — the Chinese Exclusion Act (US.07).'),
        ('African Americans', RED, 'The Compromise of 1877 ended Reconstruction’s protections (US.03); the Exodusters fled to Kansas, and Madam C.J. Walker and George Washington Carver (US.05) built enterprise under Jim Crow.'),
        ('New Immigrants & Reformers', '6A4E9C', 'Southern and Eastern European and Asian arrivals (US.07) met nativism and tenement poverty; reformers like Jane Addams at Hull-House organized mutual aid across classes.'),
    ]
    w = (12.3 - 0.3 * 3) / 4
    for i, (h, fill, body) in enumerate(cards):
        l = 0.5 + i * (w + 0.3)
        rect(s, l, 2.15, w, 4.35, fill=WHITE)
        rect(s, l, 2.15, w, 0.7, fill=fill)
        txt(s, l + 0.12, 2.2, w - 0.24, 0.6, [(h, 'Georgia', 14, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, l + 0.12, 3.0, w - 0.24, 3.4, [(body, 'Calibri', 12, INK, False)])
    txt(s, 0.5, 6.62, 12.3, 0.35, [('Verified public-domain images and named Key-Figure entries for these histories are a pending content addition.', 'Trebuchet MS', 10, LIGHT, False)])
    footer(s, 0)
    return s

# ---- build all new slides (creation order matters for reorder mapping) ----
legend = build_legend()
vocab = {c: build_vocab(c) for c in ['US.01','US.02','US.03','US.04','US.05','US.06','US.07']}
prog  = {c: build_progress(c) for c in ['US.01','US.02','US.03','US.04','US.05','US.06','US.07']}
rep = build_rep()

# ---- reorder sldIdLst ----
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)                      # in current order == prs.slides order
def idx_of(slide):                        # map slide obj -> its position in current list
    rid = slide.part.partname
    for i, s in enumerate(prs.slides):
        if s.part.partname == rid: return i
    raise KeyError
# current indices: 0..47 originals; then legend, vocab01..07, prog01..07, rep
BASE = 48
LEG = BASE
VOC = {c: BASE+1+i for i,c in enumerate(['US.01','US.02','US.03','US.04','US.05','US.06','US.07'])}
PRO = {c: BASE+8+i for i,c in enumerate(['US.01','US.02','US.03','US.04','US.05','US.06','US.07'])}
REP = BASE+15
target = [LEG,
  0, VOC['US.01'], 1,2,3,4,5,6,7,8, PRO['US.01'],
  9, VOC['US.02'], 10,11,12,13,14, PRO['US.02'],
  15, VOC['US.03'], 16,17,18,19,20,21, PRO['US.03'],
  22, VOC['US.04'], 23,24,25,26,27,28, PRO['US.04'],
  29, VOC['US.05'], 30,31,32,33,34, PRO['US.05'],
  35, VOC['US.06'], 36,37,38,39,40, PRO['US.06'],
  41, VOC['US.07'], 42,43,44,45,46,47, PRO['US.07'],
  REP]
assert sorted(target) == list(range(64)), (len(target), 'must be 0..63 once each')
new_order = [ids[i] for i in target]
for e in ids: sldIdLst.remove(e)
for e in new_order: sldIdLst.append(e)

# ---- page renumber to file position + 9pt font floor + fixed-track scan ----
FIXED = re.compile(r'DIFFERENTIATION|(?<![A-Za-z])Entry(?![A-Za-z])|Honors|EL & writing support|ON-LEVEL|BELOW-LEVEL|ABOVE-LEVEL', re.I)
flagged=[]
def all_runs(slide):
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs: yield sh, p, r
for pos, slide in enumerate(prs.slides, start=1):
    # font floor 9pt (fixes legacy 8pt footers/credits)
    for sh, p, r in all_runs(slide):
        if r.font.size is not None and r.font.size < Pt(9): r.font.size = Pt(9)
        if FIXED.search(r.text or ''): flagged.append((pos, r.text[:60]))
    # renumber: the page-number box is the small one near (12.23, 7.08)
    for sh in slide.shapes:
        if sh.has_text_frame and sh.left is not None and sh.left > Inches(11.9) and sh.top is not None and sh.top > Inches(6.9):
            t = sh.text_frame.text.strip()
            if t.isdigit():
                for p in sh.text_frame.paragraphs:
                    for r in p.runs: r.text = ''
                sh.text_frame.paragraphs[0].runs and None
                # set first run
                para = sh.text_frame.paragraphs[0]
                (para.runs[0] if para.runs else para.add_run()).text = str(pos)

# ---- alt text from on-slide captions (no file paths) ----
alt_fixed=0
for slide in prs.slides:
    caps = [sh.text_frame.text.strip() for sh in slide.shapes if sh.has_text_frame and 0 < len(sh.text_frame.text.strip()) < 140]
    cap = next((c for c in caps if any(k in c for k in ('ZOOM IN','map','photograph','Library','Archives','1878','1889','Cavalry'))), None)
    for sh in slide.shapes:
        if sh.shape_type == 13:  # picture
            cur = sh._element.find(qn('p:nvPicPr')).find(qn('p:cNvPr'))
            desc = cur.get('descr','') or ''
            if (not desc) or desc.lower().endswith(('.jpg','.jpeg','.png')) or '/' in desc or '\\' in desc:
                cur.set('descr', (cap or 'Verified public-domain historical image; see the caption and source line on this slide.'))
                alt_fixed += 1

prs.save(OUT)

# ---- patch app.xml slide count + preserve jpg content-type ----
z = zipfile.ZipFile(OUT); files = {i.filename: z.read(i.filename) for i in z.infolist()}; infos = z.infolist(); z.close()
n = len(prs.slides._sldIdLst)
if 'docProps/app.xml' in files:
    app = files['docProps/app.xml'].decode()
    app = re.sub(r'<Slides>\d+</Slides>', f'<Slides>{n}</Slides>', app)
    files['docProps/app.xml'] = app.encode()
ct = files['[Content_Types].xml'].decode()
if 'Extension="jpg"' not in ct and 'Extension="jpeg"' not in ct:
    ct = ct.replace('<Types ', '<Types ', 1)
    ct = re.sub(r'(<Types[^>]*>)', r'\1<Default Extension="jpg" ContentType="image/jpeg"/>', ct, count=1)
    files['[Content_Types].xml'] = ct.encode()
with zipfile.ZipFile(OUT, 'w', zipfile.ZIP_DEFLATED) as o:
    for i in infos: o.writestr(i, files[i.filename])

print('SAVED', OUT, '| slides:', n)
print('fixed-track flagged:', flagged or 'NONE')
print('alt-text set on pictures:', alt_fixed)
