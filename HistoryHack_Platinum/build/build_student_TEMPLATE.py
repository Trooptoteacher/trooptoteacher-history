#!/usr/bin/env python3
"""Unit 5 Student (Lean) Deck — revision pass on a WORKING COPY. Originals untouched."""
import json, zipfile, shutil, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
SRC='source_decks/Unit5_Student_Deck.pptx'; OUT='working/Unit5_Student_Deck_WORKING.pptx'
NAVY='1A2332'; NAVY2='243047'; GOLD='C9A84C'; GOLDBR='F9A825'; RED='C62828'
WHITE='FFFFFF'; SLATEL='C7D4E0'; MUTE='5C6470'; CREAM='F7F4EC'; INK='1F2430'; A250BLUE='002858'; LINEC='DAD3C2'; PANEL='2C384F'
HEAD='Georgia'; BODY='Calibri'; LABEL='Trebuchet MS'; W=13.333; H=7.5
FOOT='U.S. History Hack  ·  TroopToTeacher Technologies LLC  ·  Aligned to TN Academic Standards & TCAP EOC'
def C(h): return RGBColor.from_string(h)
STANDARDS=['US.39','US.40','US.41','US.42','US.43','US.44']

# ---- verified per-image alt + credit (from repo unit-5 catalog + teacher credit lines) ----
IMG={
 'bank-run-american-union-bank-1931':('Black-and-white photograph of a dense crowd of depositors outside the American Union Bank during a 1931 bank run.','1931 — National Archives (NARA) / Wikimedia Commons'),
 'stock-market-crash-crowd-nyse-1929':('Black-and-white photograph of a massive crowd of men filling Broad Street outside the New York Stock Exchange, October 1929.','U.S. Government, Oct 29 1929 — Wikimedia Commons'),
 'dust-bowl-farmer-fence-cimarron-1936':('Black-and-white FSA photograph by Arthur Rothstein of a farmer raising a fence post from drifted Dust Bowl sand, 1936.','Arthur Rothstein / FSA, 1936 — Library of Congress'),
 'dust-bowl-buried-machinery-dallas-sd-1936':('Black-and-white USDA photograph of farm machinery half-buried by wind-blown Dust Bowl soil, Dallas, South Dakota, 1936.','USDA, May 1936 — Wikimedia Commons'),
 'depression-soup-kitchen-breadline-nyc-1930s':('Black-and-white NARA photograph of a long line of unemployed men waiting to be fed at a New York City breadline, early 1930s.','c. 1930–1935 — National Archives (NARA) / Wikimedia Commons'),
 'migrant-mother-lange-1936':("Black-and-white FSA photograph by Dorothea Lange, 'Migrant Mother,' a worried mother with two children turned away, 1936.",'Dorothea Lange / FSA, 1936 — Library of Congress'),
 'hooverville-shantytown-1930s':('Black-and-white photograph of a Depression-era Hooverville — makeshift shacks of scrap wood and tin on the edge of a city.','Arthur Rothstein, 1936 — Library of Congress'),
 'herbert-hoover-portrait':('Formal studio portrait of President Herbert Hoover.','Underwood & Underwood, 1928 — Library of Congress'),
 'hoover-dam-construction-1930s':('Black-and-white photograph of Hoover Dam under construction in the early 1930s, canyon walls and scaffolding visible.','U.S. Bureau of Reclamation, c. 1932–1934 — National Archives'),
 'bonus-army-camp-anacostia-1932':('Black-and-white Harris & Ewing photograph of the Bonus Army camp at Anacostia Flats, Washington, D.C., 1932.','Harris & Ewing, 1932 — Library of Congress'),
 'ccc-workers-laying-stone-1933':('Black-and-white NARA photograph of Civilian Conservation Corps workers laying stone on a hillside, c. 1933.','c. 1933 — National Archives (NARA) / Wikimedia Commons'),
 'fdr-signing-social-security-1935':('Black-and-white photograph of President Franklin D. Roosevelt signing the Social Security Act as officials look on, 1935.','Acme Newspictures / Associated Press, 1935 — Library of Congress'),
 'tva-douglas-dam-construction-1942':('Color FSA/OWI photograph of Douglas Dam under construction in Tennessee, 1942, with a large crane marked US TVA.','FSA/OWI, 1942 — Library of Congress'),
 'wpa-work-pays-america-poster-1936':("Color WPA screenprint poster, 'Work Pays America! Prosperity,' with stylized figures of a farmer and worker.",'Federal Art Project / WPA, c. 1936 — Library of Congress'),
 'social-security-poster-1935':('Color 1935 federal poster promoting the new Social Security program.','Newman (artist), 1935 — Library of Congress'),
}
# orig-slide -> image slug (orig 9,10 dust-storm/okie made TEXT-ONLY; orig 34 court-packing already text-only)
SLIDE_IMG={2:'bank-run-american-union-bank-1931',3:'stock-market-crash-crowd-nyse-1929',
 8:'dust-bowl-farmer-fence-cimarron-1936',14:'depression-soup-kitchen-breadline-nyc-1930s',
 15:'migrant-mother-lange-1936',16:'hooverville-shantytown-1930s',20:'herbert-hoover-portrait',
 21:'hoover-dam-construction-1930s',22:'bonus-army-camp-anacostia-1932',26:'ccc-workers-laying-stone-1933',
 28:'tva-douglas-dam-construction-1942',
 32:'wpa-work-pays-america-poster-1936',33:'social-security-poster-1935'}
TEXTONLY={9:('US.40 · Causes: Severe Drought','A 1931–1939 drought dropped rainfall about 40% below normal, and strong winds stripped up to 75% of the topsoil, creating "black blizzards" that darkened skies as far east as Washington, D.C.'),
 10:('US.40 · Social & Economic Impacts','Displacement uprooted millions: about 3.5 million people left the Plains as "Okies," over 200,000 heading to California, as farm income fell roughly 60% and banks foreclosed on thousands of farms.'),
 27:('US.43 · Recovery & Reform Programs','The AAA and NRA raised prices and set labor codes (both later struck down by the Supreme Court), while the FDIC, SEC, the Fair Labor Standards Act, and the Social Security Act built a permanent safety net that outlasted the Depression.'),
 34:('US.44 · Court-Packing & Differing Perspectives','After the Supreme Court struck down the AAA and NRA, FDR’s 1937 plan to add up to six justices — attacked as "court packing" — failed in Congress and cost him political capital, even among allies. Yet the pressure helped shift the Court toward upholding later New Deal laws ("the switch in time that saved nine").')}

VOCAB={
 'US.39':[('Great Depression','Gran Depresión','grayt dee-PRESH-un','The worst economic downturn in U.S. history (1929–early 1940s): mass unemployment, bank failures, and widespread poverty.'),
   ('Stock Market Crash of 1929','Crash del Mercado de Valores','stok MAR-ket krash','The sudden collapse of stock prices in October 1929 ("Black Tuesday") that signaled the start of the Depression.'),
   ('Buying on Margin','Comprar con Margen','BY-ing on MAR-jin','Buying stocks with mostly borrowed money (about 10% down); the bubble burst when prices fell and loans could not be repaid.')],
 'US.40':[('Dust Bowl','Cuenca de Polvo','dust bohl','A 1930s drought-and-dust-storm disaster on the Great Plains, caused by severe drought plus over-plowing of the native grassland.'),
   ('Okies','Okies','OH-keez','A label — often used as an insult — for migrant farmers from Oklahoma, Texas, Arkansas, and Missouri who fled the Dust Bowl to California, facing discrimination and poverty.')],
 'US.41':[('Mass Unemployment','Desempleo Masivo','mass un-em-PLOY-ment','Joblessness reached about 25% (roughly 15 million people) at its 1933 peak, creating widespread poverty and desperation.'),
   ('Hoovervilles','Hoovervilles','HOO-ver-vilz','Shantytown settlements of homeless people during the Depression, named mockingly after President Herbert Hoover.')],
 'US.42':[('Rugged Individualism','Individualismo Robusto','RUG-ed in-dih-VIJ-oo-uh-liz-um','Hoover’s belief that people should succeed through their own efforts and private charity rather than direct federal aid.'),
   ('Bonus Army','Ejército del Bono','BOH-nus AR-mee','WWI veterans who marched on Washington in 1932 demanding early bonus payment; the Army dispersed them by force.'),
   ('Reconstruction Finance Corp. (RFC)','Corporación Financiera de Reconstrucción','ree-kun-STRUK-shun fy-NANS','A 1932 federal agency that lent money to banks, railroads, and businesses to prevent further collapse.')],
 'US.43':[('New Deal','New Deal','noo deel','FDR’s relief, recovery, and reform programs that expanded federal power and established the social safety net.'),
   ('Tennessee Valley Authority (TVA)','Autoridad del Valle de Tennessee','ten-eh-SEE VAL-ee aw-THOR-ih-tee','Federal agency that built dams for flood control, cheap electricity, and development in the Tennessee Valley.'),
   ('Social Security Act (1935)','Ley de Seguridad Social','SOH-shul sih-KYOOR-ih-tee','Landmark law creating old-age pensions, unemployment insurance, and aid for dependent and disabled Americans.')],
 'US.44':[('Court Packing Plan','Plan de Empaquetar la Corte','kort PAK-ing plan','FDR’s 1937 plan to add up to six Supreme Court justices after it struck down New Deal laws; seen as a threat to judicial independence.'),
   ('Socialism Accusations','Acusaciones de Socialismo','SOH-shuh-liz-um','Critics charged that the New Deal was socialist for expanding government power, though it preserved private ownership.'),
   ('FDIC','Corp. Federal de Seguro de Depósitos','FED-er-ul dee-POZ-it','New Deal agency that insures bank deposits, preventing bank runs by guaranteeing savers will not lose their money.')],
}
# CFU DOK-2 items: original options with correct at index 1 (longest). PERM re-orders to de-bias position.
PCHECK={
 'US.39':('What role did "buying on margin" and overextension of credit play in creating economic instability?',
   ['Margin buying had no connection to the crash.','Margin buying inflated stock prices artificially, and when prices fell, margin calls forced panic selling that accelerated the collapse while widespread consumer debt created defaults.','Credit was carefully regulated in the 1920s.','Only wealthy investors bought on margin.']),
 'US.40':('What combination of human farming practices and natural weather patterns created the Dust Bowl?',
   ['The Dust Bowl was entirely caused by natural drought.','Farmers plowed up native grasses, practiced monoculture without conservation, and severe drought combined with strong winds eroded the exposed topsoil.','Modern farming techniques prevented soil erosion.','The Dust Bowl only affected a small local area.']),
 'US.41':('How did mass unemployment during the Great Depression affect individuals, families, and communities?',
   ['Unemployment had limited social effects.','With 25% unemployment, families faced psychological trauma, role reversals, increased divorce, child labor, and relied on mutual aid while communities experienced population loss and service decline.','Government programs immediately solved unemployment.','Only industrial workers experienced unemployment.']),
 'US.42':('How did President Hoover’s philosophy of "Rugged Individualism" shape his response to the Great Depression?',
   ['It led to aggressive government intervention.','His belief that individuals should solve problems without government help led to limited federal response, reliance on voluntary cooperation, and opposition to direct relief programs.','He immediately provided direct payments to unemployed workers.','Rugged individualism had no impact on his policies.']),
 'US.43':('How did the "Three R’s" (Relief, Recovery, Reform) organize Roosevelt’s approach to ending the Great Depression?',
   ['FDR used only one approach throughout the New Deal.','Relief provided immediate help to unemployed (CCC, WPA), Recovery programs restored economic activity (AAA, NRA), and Reform programs fixed underlying problems (FDIC, SEC, Social Security).','The Three R’s were unrelated programs.','Roosevelt opposed the Three R’s approach.']),
 'US.44':('What were the main economic effects of New Deal policies on employment, banking, and government spending?',
   ['New Deal policies had no measurable economic impact.','Unemployment fell from 25% to 14%, bank runs ended with FDIC, manufacturing increased 60%, and government spending grew from 3% to 10% of GDP.','The economy worsened under New Deal policies.','Only banking was affected by New Deal programs.']),
}
# de-bias permutations (source-index order); correct (orig idx1) lands at varied positions C,A,D,B,A,C
PERM={'US.39':[0,2,1,3],'US.40':[1,0,2,3],'US.41':[0,2,3,1],'US.42':[0,1,2,3],'US.43':[1,0,2,3],'US.44':[0,2,1,3]}
def reordered(std):
    q,opts=PCHECK[std]; p=PERM[std]
    return q,[opts[i] for i in p], p.index(1)  # (stem, options, correct_index)

# ---------- helpers ----------
def rect(s,x,y,w,h,fill,line=None,lw=1.0,rounded=False):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.shadow.inherit=False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=C(fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=C(line); sh.line.width=Pt(lw)
    return sh
def txt(s,x,y,w,h,runs,size,color,font=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.0,fit=False):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    if fit: tf.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if isinstance(runs,str): runs=[[(runs,{})]]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if sp: p.line_spacing=sp
        for t,opt in para:
            r=p.add_run(); r.text=t; f=r.font
            f.size=Pt(opt.get('size',size)); f.name=opt.get('font',font); f.bold=opt.get('bold',bold)
            f.italic=opt.get('italic',italic); f.color.rgb=C(opt.get('color',color))
    return tb
def footer(s,dark):
    txt(s,0.5,7.08,11.0,0.3,FOOT,9,SLATEL if dark else MUTE,font=LABEL)
def new_slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[0])
    for ph in list(s.placeholders): ph._element.getparent().remove(ph._element)
    return s
def set_alt(pic,alt): pic._element.xpath('.//p:cNvPr')[0].set('descr',alt)

TITLE={'US.39':'Causes of the Great Depression','US.40':'The Dust Bowl','US.41':'Impact of the Great Depression on the American People','US.42':"Herbert Hoover's Response",'US.43':"FDR's New Deal Programs",'US.44':'Effects & Controversies of New Deal Policies'}

shutil.copy(SRC,OUT); prs=Presentation(OUT)
prs.core_properties.subject='Unit 5 — The Great Depression & the New Deal (US.39–US.44)'
prs.core_properties.title='U.S. History Hack — Unit 5 Course Standard Student Deck'
prs.core_properties.author='TroopToTeacher Technologies LLC'
sl=list(prs.slides)

# (A) image slides: alt + visible credit
for n,slug in SLIDE_IMG.items():
    s=sl[n-1]; alt,credit=IMG[slug]
    for sh in s.shapes:
        if sh.shape_type==13: set_alt(sh,alt)
    txt(s,7.35,6.8,5.45,0.26,[[('Source: ',{'color':MUTE}),(credit,{'color':A250BLUE,'italic':True})]],9,MUTE,font=BODY,align=PP_ALIGN.RIGHT)

# (B) de-dup KEY POINT on orig slide 4
for sh in sl[3].shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith('Structural weakness deepened the crisis') and abs(Emu(sh.top).inches-3.15)<0.3:
        tf=sh.text_frame
        for p in tf.paragraphs:
            for r in p.runs: r.text=''
        (tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()).text=('Even before the 1929 crash, the boom rested on shaky supports — too many goods, too much debt, and too little regulation — so joblessness climbed from 3.2% to 23.6%.')

# (C) text-only slides (orig 9,10,34): remove picture + orphan ZOOM IN; add KEY POINT card
for n,(cap,body) in TEXTONLY.items():
    s=sl[n-1]
    for sh in list(s.shapes):
        if sh.shape_type==13: sh._element.getparent().remove(sh._element)
        elif sh.has_text_frame and sh.text_frame.text.strip()=='ZOOM IN': sh._element.getparent().remove(sh._element)
    rect(s,1.6,2.6,10.13,3.7,WHITE,line=LINEC,lw=1.0); rect(s,1.6,2.6,10.13,0.55,NAVY)
    txt(s,1.9,2.6,9.5,0.55,'KEY POINT',13,GOLDBR,font=LABEL,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,2.1,3.35,9.13,2.7,body,16,INK,font=BODY,anchor=MSO_ANCHOR.TOP,sp=1.12,fit=True)

# (D) DI headline auto-fit (fixes clipped long headlines in render)
for i,s in enumerate(sl):
    for sh in s.shapes:
        if sh.has_text_frame and abs(Emu(sh.top).inches-0.85)<0.15 and Emu(sh.width).inches>11 and sh.text_frame.text.strip():
            sh.text_frame.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

# (E) new slides: legend + per-standard vocab & progress check
def build_legend(prs):
    s=new_slide(prs); rect(s,0,0,W,H,NAVY); rect(s,0,0,W,0.16,GOLD)
    txt(s,0.7,0.55,12,0.4,'UNIVERSAL DESIGN FOR LEARNING · MTSS',12,GOLDBR,font=LABEL,bold=True)
    txt(s,0.7,0.95,12,0.8,'How This Unit Supports Every Learner',32,WHITE,font=HEAD,bold=True)
    labels=[('CORE PATH','The grade-level lesson every student does — same TN standard, same rigor, same ceiling.'),
      ('SUPPORT OPTION','An optional scaffold that lowers a barrier, not the goal. Available to anyone.'),
      ('LANGUAGE SUPPORT','Vocabulary, pronunciation, Spanish, and sentence frames for multilingual learners — open to all.'),
      ('RESPONSE CHOICE','Show what you know your way: write it · say & record it · diagram & label it. Same learning target.'),
      ('PROGRESS CHECK','A quick DOK-2 check and confidence rating so you and your teacher see the next step.'),
      ('EXTENSION','An optional deeper challenge (higher DOK). The standard stays the same.')]
    x0,y0,cw,ch,gx,gy=0.7,2.0,5.95,1.15,0.25,0.18
    for i,(lab,desc) in enumerate(labels):
        x=x0+(i%2)*(cw+gx); y=y0+(i//2)*(ch+gy)
        rect(s,x,y,cw,ch,NAVY2,line=GOLD,lw=0.75,rounded=True)
        txt(s,x+0.25,y+0.12,cw-0.5,0.3,lab,13,GOLDBR,font=LABEL,bold=True)
        txt(s,x+0.25,y+0.46,cw-0.5,ch-0.55,desc,11.5,SLATEL,font=BODY,sp=1.03)
    rect(s,0.7,5.95,12.0,0.95,'11202F',line=GOLDBR,lw=1.0,rounded=True)
    txt(s,0.95,6.06,11.5,0.78,[[('Options are available by design and support access without lowering the learning goal.  ',{'bold':True,'color':WHITE}),('Supports work alongside, never in place of, required IEP or 504 accommodations.',{'color':GOLDBR,'bold':True})]],13,WHITE,font=BODY,anchor=MSO_ANCHOR.MIDDLE,sp=1.05)
    footer(s,True); return s
def build_vocab(prs,std):
    s=new_slide(prs); rect(s,0,0,W,H,CREAM); rect(s,0,0,W,0.16,NAVY)
    txt(s,0.6,0.42,10,0.3,f'{std} · KEY VOCABULARY / LANGUAGE SUPPORT',12,RED,font=LABEL,bold=True)
    txt(s,0.6,0.78,12.1,0.7,'Key Vocabulary — Language Support',30,NAVY,font=HEAD,bold=True)
    terms=VOCAB[std]; n=len(terms); cw=3.9 if n==3 else 5.4; gap=0.3; total=n*cw+(n-1)*gap; x0=(W-total)/2; y=1.95; ch=4.15; hh=1.05
    for i,(en,es,say,defn) in enumerate(terms):
        x=x0+i*(cw+gap)
        rect(s,x,y,cw,ch,WHITE,line=LINEC,lw=1.0); rect(s,x,y,cw,hh,NAVY)
        txt(s,x+0.18,y+0.08,cw-0.36,0.56,en,14.5,WHITE,font=BODY,bold=True,sp=0.98)
        txt(s,x+0.18,y+0.62,cw-0.36,0.22,es,10.5,GOLDBR,font=BODY,italic=True)
        txt(s,x+0.18,y+0.83,cw-0.36,0.2,'say:  '+say,10,SLATEL,font=BODY,italic=True)
        txt(s,x+0.18,y+hh+0.12,cw-0.36,ch-hh-0.24,defn,12,INK,font=BODY,sp=1.06,fit=True)
    txt(s,0.6,6.45,12.1,0.45,[[('LANGUAGE SUPPORT  ',{'bold':True,'color':RED,'font':LABEL}),('Spanish term, a "say:" pronunciation guide, and a plain-language definition support multilingual learners and are available to everyone. Full bilingual Frayer models are in the unit vocabulary set.',{'color':MUTE,'italic':True})]],10,MUTE,font=BODY,sp=1.02)
    footer(s,False); return s
def build_pcheck(prs,std):
    stem,opts,ci=reordered(std)
    s=new_slide(prs); rect(s,0,0,W,H,NAVY); rect(s,0,0,W,0.16,GOLD)
    txt(s,0.6,0.48,9,0.3,f'{std} · PROGRESS CHECK',12,GOLDBR,font=LABEL,bold=True)
    txt(s,0.6,0.82,8.5,0.6,'Progress Check',28,WHITE,font=HEAD,bold=True)
    rect(s,9.9,0.9,2.8,0.5,PANEL,line=GOLD,lw=0.75,rounded=True)
    txt(s,9.9,0.9,2.8,0.5,'DOK 2 · SELF-CHECK',11,GOLDBR,font=LABEL,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    rect(s,0.6,1.62,12.1,1.05,NAVY2,line=SLATEL,lw=0.5,rounded=True)
    txt(s,0.9,1.72,11.5,0.9,stem,16,WHITE,font=BODY,bold=True,anchor=MSO_ANCHOR.MIDDLE,sp=1.03)
    for i,o in enumerate(opts):
        oy=2.92+i*0.72
        rect(s,0.6,oy,12.1,0.62,PANEL,line='3A475E',lw=1.0,rounded=True)
        rect(s,0.78,oy+0.12,0.38,0.38,NAVY); txt(s,0.78,oy+0.12,0.38,0.38,'ABCD'[i],13,SLATEL,font=LABEL,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,1.33,oy,10.9,0.62,o,12.5,WHITE,font=BODY,anchor=MSO_ANCHOR.MIDDLE,sp=1.0,fit=True)
    rect(s,0.6,5.95,12.1,0.95,NAVY2,line=GOLDBR,lw=1.0,rounded=True)
    txt(s,0.9,6.03,11.6,0.8,[[('Commit your answer first',{'bold':True,'color':GOLDBR}),(' — then discuss. Your teacher has the answer key.   How sure are you?  ',{'color':SLATEL}),('I’ve got this  ·  Mostly  ·  Still shaky  ·  Lost',{'bold':True,'color':WHITE})]],12,SLATEL,font=BODY,anchor=MSO_ANCHOR.MIDDLE,sp=1.04)
    footer(s,True); return s

def build_rep(prs):
    s=new_slide(prs); rect(s,0,0,W,H,CREAM); rect(s,0,0,W,0.16,RED)
    txt(s,0.6,0.42,12,0.3,'US.41 & US.43 · WHOSE STORY THE NEW DEAL REACHED',12,RED,font=LABEL,bold=True)
    txt(s,0.6,0.78,12.1,0.7,'Reached — and Left Out',30,NAVY,font=HEAD,bold=True)
    blocks=[('Black Americans','The AAA\u2019s crop-reduction payments went to landowners, pushing many Black sharecroppers and tenant farmers off the land, and Social Security (1935) at first excluded farm and domestic workers \u2014 jobs held by most Black workers. Yet FDR\u2019s informal \u201cBlack Cabinet,\u201d including Mary McLeod Bethune and Robert Weaver, advised the administration, and Black voters shifted toward the Democratic Party.'),
      ('Mexican Americans','Amid Depression job competition, \u201cMexican Repatriation\u201d (about 1929\u20131936) pressured an estimated 400,000 to more than one million people of Mexican descent \u2014 many of them U.S.-born citizens \u2014 to leave for Mexico. Mexican and Filipino farmworkers in California often faced lower pay and worse conditions than white \u201cOkie\u201d migrants.'),
      ('Indigenous Peoples','The Indian Reorganization Act of 1934 (the \u201cIndian New Deal\u201d), led by John Collier, ended the Dawes-era allotment policy, restored some tribal self-government and land, and supported Native culture \u2014 a major reversal, though its results were mixed and debated. A CCC Indian Division employed thousands.')]
    cw=3.95; gap=0.2; x0=0.6; y=1.95; ch=4.2; hh=0.5
    for i,(h,b) in enumerate(blocks):
        x=x0+i*(cw+gap)
        rect(s,x,y,cw,ch,WHITE,line=LINEC,lw=1.0); rect(s,x,y,cw,hh,NAVY)
        txt(s,x+0.18,y,cw-0.36,hh,h,14,WHITE,font=BODY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+0.18,y+hh+0.14,cw-0.36,ch-hh-0.28,b,11.5,INK,font=BODY,sp=1.06,fit=True)
    txt(s,0.6,6.5,12.1,0.4,[[('History Hack-authored instructional synthesis grounded in the standard historical record \u2014 not a primary source. ',{'italic':True,'color':MUTE}),('Verified public-domain images and primary sources for these histories are a pending content addition.',{'italic':True,'color':RED})]],9.5,MUTE,font=BODY,sp=1.0)
    footer(s,False); return s
legend=build_legend(prs); vocab={s:build_vocab(prs,s) for s in STANDARDS}; pcheck={s:build_pcheck(prs,s) for s in STANDARDS}; rep=build_rep(prs)
order=[legend]
for bi,std in enumerate(STANDARDS):
    b=bi*6; order+=[sl[b],vocab[std],sl[b+1],sl[b+2],sl[b+3],sl[b+4],sl[b+5],pcheck[std]]
order.append(rep)
cur=list(prs.slides); lst=prs.slides._sldIdLst; kids=list(lst); s2el={id(cur[i]):kids[i] for i in range(len(cur))}
for el in kids: lst.remove(el)
for s in order: lst.append(s2el[id(s)])

# (F) font floor 9pt + blanket auto-fit (shrinks only overflowing frames) across all slides
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is not None and r.font.size < Pt(9): r.font.size=Pt(9)
            if sh.text_frame.text.strip() and not (Emu(sh.left or 0).inches>11.9 and Emu(sh.top or 0).inches>6.9):
                sh.text_frame.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

# (G) page numbers -> file position (1..N); add where missing (new slides)
for i,s in enumerate(prs.slides,1):
    num=None
    for sh in s.shapes:
        if sh.has_text_frame and Emu(sh.left or 0).inches>11.9 and Emu(sh.top or 0).inches>6.9 and sh.text_frame.text.strip().isdigit():
            num=sh; break
    if num is not None:
        r=num.text_frame.paragraphs[0].runs
        (r[0] if r else num.text_frame.paragraphs[0].add_run()).text=str(i)
    else:
        # detect dark bg (navy new slides) vs cream (vocab)
        dark=any(getattr(sh,'fill',None) and sh.fill.type==1 and str(sh.fill.fore_color.rgb)==NAVY and Emu(sh.width or 0).inches>13 for sh in s.shapes)
        txt(s,12.23,7.08,0.6,0.3,str(i),9,SLATEL if dark else MUTE,font=LABEL,align=PP_ALIGN.RIGHT)

prs.save(OUT)
# (H) packaging: jpg content-type + app.xml slide count
z=zipfile.ZipFile(OUT); files={it.filename:z.read(it.filename) for it in z.infolist()}; infos=z.infolist(); z.close()
ct=files['[Content_Types].xml'].decode()
if 'Extension="jpg"' not in ct:
    ct=ct.replace('<Default Extension="jpeg" ContentType="image/jpeg"/>','<Default Extension="jpeg" ContentType="image/jpeg"/><Default Extension="jpg" ContentType="image/jpeg"/>')
    files['[Content_Types].xml']=ct.encode()
if 'docProps/app.xml' in files:
    app=files['docProps/app.xml'].decode()
    app=re.sub(r'<Slides>\d+</Slides>','<Slides>50</Slides>',app)
    files['docProps/app.xml']=app.encode()
with zipfile.ZipFile(OUT,'w',zipfile.ZIP_DEFLATED) as o:
    for it in infos: o.writestr(it, files[it.filename])
print('SAVED',OUT,'slides:',len(list(Presentation(OUT).slides)),'| debias correct letters:',{s:'ABCD'[reordered(s)[2]] for s in STANDARDS})
