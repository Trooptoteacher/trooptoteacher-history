# -*- coding: utf-8 -*-
"""Insert full-color projection map slides for the two color-dependent maps
(US.01 railroad map, US.06 foreign-born map) into the student + teacher decks,
right after each standard's source slide. Honors the workbook's on-page note
that a full-color version lives on the projection slide. Deck palette matched."""
import json, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

IMG = json.load(open('analysis/unit1_images.json'))
NAVY='1A2332'; GOLD='C9A84C'; LIGHT='C7D4E0'; WHITE='FFFFFF'; INK='1F2430'
def C(h): return RGBColor.from_string(h)

def blank_layout(prs):
    best=None
    for lay in prs.slide_layouts:
        n=len(lay.placeholders)
        if best is None or n < len(best.placeholders): best=lay
    return best

def set_bg(slide, hexc=NAVY):
    cSld = slide._element.find(qn('p:cSld'))
    bg = cSld.find(qn('p:bg'))
    if bg is not None: cSld.remove(bg)
    from pptx.oxml import parse_xml
    xml=('<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
         'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
         f'<p:bgPr><a:solidFill><a:srgbClr val="{hexc}"/></a:solidFill>'
         '<a:effectLst/></p:bgPr></p:bg>')
    cSld.insert(0, parse_xml(xml))

def textbox(slide, l,t,w,h, text, size, color, bold=True, align=PP_ALIGN.LEFT, italic=False):
    tb=slide.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=MSO_ANCHOR.TOP
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=C(color); f.name='Calibri'
    return tb

def add_map_slide(prs, rec, code, after_1based):
    SW,SH = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(blank_layout(prs))
    # strip any inherited placeholders
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    set_bg(slide, NAVY)
    m=Inches(0.55)
    # gold eyebrow + title
    textbox(slide, m, Inches(0.30), SW-2*m, Inches(0.4),
            f"GEOGRAPHER'S LENS  ·  {code}  —  COLOR MAP FOR PROJECTION", 15, GOLD, True, PP_ALIGN.LEFT)
    textbox(slide, m, Inches(0.66), SW-2*m, Inches(0.55), rec['title'], 20, WHITE, True, PP_ALIGN.LEFT)
    # image scaled to fit the middle band
    iw,ih = Image.open(rec['file']).size
    ar = iw/ih
    avail_w = SW-2*m
    avail_h = SH - Inches(0.30+0.66+0.55+1.15)  # leave room for header + caption band
    w = avail_w; h = int(w/ar)
    if h > avail_h: h=avail_h; w=int(h*ar)
    left = int((SW - w)/2); top = Inches(1.45)
    slide.shapes.add_picture(rec['file'], left, top, width=w, height=h)
    # caption + citation band
    cap_top = top + h + Inches(0.08)
    textbox(slide, m, cap_top, SW-2*m, Inches(0.9),
            f"This is the full-color version of the black-and-white map on the workbook Geographer's Lens page for {code}.  "
            f"Source: {rec['citation']}  {rec.get('rights','')}", 10, LIGHT, False, PP_ALIGN.LEFT, italic=True)
    # NOTE: slides are APPENDED at the end (projection appendix) so no existing
    # slide number shifts — this keeps the workbook's Cornell/deck references valid.

def process(path, codes):
    prs=Presentation(path)
    n0=len(list(prs.slides))
    for code in codes:            # appended at end, in order
        add_map_slide(prs, IMG[code]['map'], code, None)
    prs.save(path)
    n1=len(list(prs.slides))
    print(f'  {path} -> {n1} slides (was {n0}; +{n1-n0} projection maps appended at end: {", ".join(codes)})')

print('STUDENT deck:')
process('deliverables/Unit1_Student_Deck_CourseStandard.pptx', ['US.01','US.06'])
print('TEACHER deck:')
process('deliverables/Unit1_Teacher_Deck_CourseStandard.pptx', ['US.01','US.06'])
print('done')
