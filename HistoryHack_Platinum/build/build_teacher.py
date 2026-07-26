# -*- coding: utf-8 -*-
"""Unit 1 Teacher (Full) Deck — PLATINUM edit-in-place.
Adds UDL/MTSS legend + Representation; de-biases CFU + ANSWER REVEAL IN SYNC with
the student progress checks (moves correct text, the gold highlight bar, the ✓, and
the 'Correct answer:' line together); adds TEACHER DECISION BAND + SCAFFOLD FADING to
every wrap-up; removes fixed-track labels in bodies AND notes; alt text; 9pt floor;
renumber. Real images intact. Originals untouched."""
import re, zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import unit1_data as D

IN  = 'working/Unit1_Teacher_Deck_WORKING.pptx'
OUT = 'deliverables/Unit1_Teacher_Deck_CourseStandard.pptx'
NAVY='1A2332'; GOLD='C9A84C'; LIGHT='C7D4E0'; WHITE='FFFFFF'; INK='1F2430'
GREEN='1B5E20'; RED='C62828'; BLUE='002858'; GREY='5C6470'; ROWDK='2C384F'
FOOT=D.UNIT['footer']
def C(h): return RGBColor.from_string(h)
prs=Presentation(IN); BLANK=prs.slide_layouts[0]
STDS=['US.01','US.02','US.03','US.04','US.05','US.06','US.07']

def set_bg(s,hexc=NAVY):
    cSld=s._element.find(qn('p:cSld'))
    for ex in cSld.findall(qn('p:bg')): cSld.remove(ex)
    bg=cSld.makeelement(qn('p:bg'),{}); pr=bg.makeelement(qn('p:bgPr'),{})
    fill=pr.makeelement(qn('a:solidFill'),{}); clr=fill.makeelement(qn('a:srgbClr'),{'val':hexc})
    fill.append(clr); pr.append(fill); pr.append(pr.makeelement(qn('a:effectLst'),{})); bg.append(pr); cSld.insert(0,bg)
def new_slide():
    s=prs.slides.add_slide(BLANK)
    for sh in list(s.shapes): sh._element.getparent().remove(sh._element)
    set_bg(s); return s
def rect(s,l,t,w,h,fill=None,line=None,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=C(line); sp.line.width=Pt(lw)
    return sp
def txt(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=2):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in ('left','right','top','bottom'): setattr(tf,'margin_'+m,Emu(0))
    if isinstance(runs[0],tuple): runs=[runs]
    for pi,para in enumerate(runs):
        p=tf.paragraphs[0] if pi==0 else tf.add_paragraph(); p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0)
        for (t_,fn,sz,col,b) in para:
            r=p.add_run(); r.text=t_; r.font.name=fn; r.font.size=Pt(sz); r.font.color.rgb=C(col); r.font.bold=b
    return tb
def chip(s,l,t,w,label,fill,fg=WHITE,sz=11):
    rect(s,l,t,w,0.4,fill=fill); txt(s,l,t,w,0.4,[(label,'Trebuchet MS',sz,fg,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def footer(s,page):
    txt(s,0.5,7.08,11.0,0.3,[(FOOT,'Trebuchet MS',9,LIGHT,False)])
    txt(s,12.23,7.08,0.6,0.3,[(str(page),'Trebuchet MS',9,LIGHT,False)],align=PP_ALIGN.RIGHT)

def build_legend():
    s=new_slide()
    txt(s,0.5,0.45,12.3,0.7,[('How This Deck Works — UDL & MTSS','Georgia',30,WHITE,True)])
    txt(s,0.5,1.18,12.3,0.5,[('Universal supports built into every standard. One firm learning goal; flexible ways to reach it.','Georgia',15,LIGHT,False)])
    items=[('CORE PATH',GOLD,NAVY,'The essential instruction every student receives.'),
           ('SUPPORT OPTION',GREEN,WHITE,'An optional scaffold that keeps the goal, not lowers it.'),
           ('LANGUAGE SUPPORT',BLUE,WHITE,'Vocabulary, pronunciation, and Spanish cognates for access.'),
           ('RESPONSE CHOICE',RED,WHITE,'Show learning by writing, saying/recording, or diagramming.'),
           ('PROGRESS CHECK','6A4E9C',WHITE,'A quick DOK-2/3 check to guide reteach or extend.'),
           ('EXTENSION','2F6F4E',WHITE,'A deeper challenge once the goal is met.')]
    x0,y0,w,h,gx,gy=0.5,1.95,3.95,1.15,0.22,0.28
    for i,(lab,fill,fg,g) in enumerate(items):
        col,row=i%3,i//3; l=x0+col*(w+gx); t=y0+row*(h+gy)
        rect(s,l,t,w,h,fill=WHITE); chip(s,l,t,w,lab,fill,fg=fg,sz=12); txt(s,l+0.14,t+0.5,w-0.28,h-0.55,[(g,'Calibri',12,INK,False)])
    rect(s,0.5,5.95,12.3,0.95,fill='232F40',line=GOLD,lw=1.5)
    txt(s,0.7,6.08,11.9,0.8,[[('Options are available by design and support access without lowering the learning goal.  ','Calibri',12.5,WHITE,True)],
                            [('Supports work alongside — never in place of — required IEP or 504 accommodations.','Calibri',12.5,LIGHT,False)]])
    footer(s,0); return s

def build_rep():
    s=new_slide()
    chip(s,0.5,0.3,3.9,'REACHED — AND LEFT OUT',GOLD,fg=NAVY,sz=11)
    txt(s,0.5,0.9,12.3,0.65,[('Whose story the industrial era told — and who paid its costs','Georgia',22,WHITE,True)])
    txt(s,0.5,1.55,12.3,0.5,[('History Hack-authored synthesis grounded in this unit’s sourced record — read across the standards, not around them.','Georgia',13,LIGHT,False)])
    cards=[('American Indian Nations',GREEN,'Westward expansion and the railroad (US.01) crossed treaty lands; reservation, assimilation, boarding-school, and Dawes Act policy (US.02) cut tribal holdings from 138 to 48 million acres by 1934.'),
           ('Chinese Laborers',BLUE,'Chinese workers built the most dangerous stretches of the Central Pacific (US.01) yet faced wage discrimination and, in 1882, the first federal law to bar a nationality — the Chinese Exclusion Act (US.07).'),
           ('African Americans',RED,'The Compromise of 1877 ended Reconstruction’s protections (US.03); the Exodusters fled to Kansas, and Madam C.J. Walker and George Washington Carver (US.05) built enterprise under Jim Crow.'),
           ('New Immigrants & Reformers','6A4E9C','Southern and Eastern European and Asian arrivals (US.07) met nativism and tenement poverty; reformers like Jane Addams at Hull-House organized mutual aid across classes.')]
    w=(12.3-0.3*3)/4
    for i,(h,fill,body) in enumerate(cards):
        l=0.5+i*(w+0.3); rect(s,l,2.15,w,4.35,fill=WHITE); rect(s,l,2.15,w,0.7,fill=fill)
        txt(s,l+0.12,2.2,w-0.24,0.6,[(h,'Georgia',14,WHITE,True)],anchor=MSO_ANCHOR.MIDDLE); txt(s,l+0.12,3.0,w-0.24,3.4,[(body,'Calibri',12,INK,False)])
    txt(s,0.5,6.62,12.3,0.35,[('Verified public-domain images and named Key-Figure entries for these histories are a pending content addition.','Trebuchet MS',10,LIGHT,False)])
    footer(s,0); return s

# ---------- helpers to edit existing slides ----------
def set_shape_text(sh,newtext):
    tf=sh.text_frame; p=tf.paragraphs[0]
    if not p.runs: p.add_run()
    p.runs[0].text=newtext
    for r in p.runs[1:]: r.text=''
    for extra in tf.paragraphs[1:]:
        for r in extra.runs: r.text=''

def opt_text_shapes(slide):
    got=[sh for sh in slide.shapes if sh.has_text_frame and sh.left is not None
         and Inches(1.28)<sh.left<Inches(1.45) and sh.width is not None and sh.width>Inches(9)]
    return sorted(got,key=lambda s:s.top)

def rowbars(slide):
    got=[]
    for sh in slide.shapes:
        try: fh=str(sh.fill.fore_color.rgb) if sh.fill.type==1 else ''
        except Exception: fh=''
        if sh.left is not None and Inches(0.5)<sh.left<Inches(0.7) and sh.width and sh.width>Inches(11.8) \
           and sh.height and Inches(0.55)<sh.height<Inches(0.7) and fh in (GOLD,ROWDK):
            got.append(sh)
    return sorted(got,key=lambda s:s.top)

def debias(slide,code,is_reveal):
    opts,tgt=D.debiased_options(code); order=['A','B','C','D']
    ots=opt_text_shapes(slide)
    if len(ots)!=4: raise RuntimeError(f'{code} reveal={is_reveal}: found {len(ots)} option shapes')
    for sh,L in zip(ots,order): set_shape_text(sh,opts[L])
    if is_reveal:
        bars=rowbars(slide); ti=order.index(tgt)
        if len(bars)==4:
            for k,b in enumerate(bars):
                b.fill.solid(); b.fill.fore_color.rgb=C(GOLD if k==ti else ROWDK)
            for sh in slide.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip()=='✓' and sh.left is not None and sh.left>Inches(11.5):
                    sh.top=bars[ti].top; sh.height=bars[ti].height
        for sh in slide.shapes:
            if sh.has_text_frame and 'Correct answer' in sh.text_frame.text:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text: r.text=re.sub(r'(Correct answer:\s*)[A-D]', r'\g<1>'+tgt, r.text)

def add_decision_band(slide):
    rect(slide,0.5,6.42,12.33,0.56,fill='232F40',line=GOLD,lw=1.0)
    txt(slide,0.66,6.46,12.0,0.5,[
        [('TEACHER DECISION BAND   ','Trebuchet MS',10,GOLD,True),
         ('PROGRESS CHECK → identify barrier → select support → reteach → recheck → extend','Calibri',10,WHITE,False)],
        [('SCAFFOLD FADING   ','Trebuchet MS',10,GOLD,True),
         ('Guided → Light → independent as evidence grows; required IEP/504 supports do not fade unless the authorized team changes them.','Calibri',10,LIGHT,False)],
    ],sp_after=0)

# ---- label replacements (bodies + notes) ----
REPL=[('(EL & writing support)','· LANGUAGE SUPPORT'),('HONORS / EXTENSION','EXTENSION'),
      ('HONORS/EXTENSION','EXTENSION'),('DIFFERENTIATION','SUPPORT OPTIONS'),
      ('Entry:','SUPPORT OPTION:'),('Honors:','EXTENSION:'),('HONORS','EXTENSION')]
def relabel_runs(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            if not r.text: continue
            s=r.text
            for a,b in REPL: s=s.replace(a,b)
            if s!=r.text: r.text=s

# ---- CFU/reveal indices ----
CFU_IDX ={'US.01':14,'US.02':28,'US.03':44,'US.04':59,'US.05':74,'US.06':89,'US.07':105}
REV_IDX ={'US.01':15,'US.02':29,'US.03':45,'US.04':60,'US.05':75,'US.06':90,'US.07':106}
WRAP_IDX=[16,30,46,61,76,91,107]

for code,i in CFU_IDX.items(): debias(prs.slides[i],code,False)
for code,i in REV_IDX.items(): debias(prs.slides[i],code,True)
for i in WRAP_IDX: add_decision_band(prs.slides[i])

# global relabel (bodies + notes) + 9pt floor + alt text
alt_fixed=0
for slide in prs.slides:
    for sh in slide.shapes:
        if sh.has_text_frame: relabel_runs(sh.text_frame)
    if slide.has_notes_slide: relabel_runs(slide.notes_slide.notes_text_frame)
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is not None and r.font.size<Pt(9): r.font.size=Pt(9)
        if sh.shape_type==13:
            cur=sh._element.find(qn('p:nvPicPr')).find(qn('p:cNvPr')); desc=cur.get('descr','') or ''
            caps=[x.text_frame.text.strip() for x in slide.shapes if x.has_text_frame and 0<len(x.text_frame.text.strip())<130]
            cap=next((c for c in caps if any(k in c for k in ('Library','Archives','Congress','photograph','map','18','Cavalry','NARA','LOC'))),None)
            if (not desc) or desc.lower().endswith(('.jpg','.jpeg','.png')) or '/' in desc or '\\' in desc:
                cur.set('descr',cap or 'Verified public-domain historical image; see the caption and source line on this slide.'); alt_fixed+=1

# ---- insert legend (after front matter, pos 3) and representation (before close) ----
legend=build_legend(); rep=build_rep()
sldIdLst=prs.slides._sldIdLst; ids=list(sldIdLst)  # 0..108 orig, 109 legend, 110 rep
target=[0,1,2,109]+list(range(3,108))+[110,108]
assert sorted(target)==list(range(111)),(len(target))
order=[ids[i] for i in target]
for e in ids: sldIdLst.remove(e)
for e in order: sldIdLst.append(e)

# ---- renumber page-number boxes to file position ----
for pos,slide in enumerate(prs.slides,start=1):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.left is not None and sh.left>Inches(11.9) and sh.top is not None and sh.top>Inches(6.9):
            if sh.text_frame.text.strip().isdigit():
                para=sh.text_frame.paragraphs[0]
                (para.runs[0] if para.runs else para.add_run()).text=str(pos)
                for r in para.runs[1:]: r.text=''

prs.save(OUT)
z=zipfile.ZipFile(OUT); files={i.filename:z.read(i.filename) for i in z.infolist()}; infos=z.infolist(); z.close()
n=len(prs.slides._sldIdLst)
if 'docProps/app.xml' in files:
    app=files['docProps/app.xml'].decode(); app=re.sub(r'<Slides>\d+</Slides>',f'<Slides>{n}</Slides>',app); files['docProps/app.xml']=app.encode()
ct=files['[Content_Types].xml'].decode()
if 'Extension="jpg"' not in ct and 'Extension="jpeg"' not in ct:
    ct=re.sub(r'(<Types[^>]*>)',r'\1<Default Extension="jpg" ContentType="image/jpeg"/>',ct,count=1); files['[Content_Types].xml']=ct.encode()
with zipfile.ZipFile(OUT,'w',zipfile.ZIP_DEFLATED) as o:
    for i in infos: o.writestr(i,files[i.filename])

# verify de-bias keys
print('SAVED',OUT,'| slides:',n,'| alt set:',alt_fixed)
for code in STDS:
    _,k=D.debiased_options(code); print(f'  {code} de-biased key -> {k}')
