#!/usr/bin/env python3
"""Derive the exact Unit 5 slide crosswalk from the FINAL edited decks.
Every slide number is read from the PPTX — nothing is estimated."""
import json, re
from pptx import Presentation
STDS=['US.39','US.40','US.41','US.42','US.43','US.44']

def slide_text(s):
    out=[]
    for sh in s.shapes:
        if sh.has_text_frame: out.append(sh.text_frame.text)
    return '\n'.join(out)

def classify_student(txt):
    if 'How This Unit Supports' in txt: return ('','legend')
    std=next((s for s in STDS if s in txt),'')
    if 'KEY VOCABULARY / LANGUAGE SUPPORT' in txt: t='vocab'
    elif 'PROGRESS CHECK' in txt: t='progress'
    elif 'SOURCE IT FIRST' in txt: t='source'
    elif 'THREE PERSPECTIVES' in txt: t='perspectives'
    elif 'DIRECT INSTRUCTION' in txt: t='di'
    elif 'TN STANDARD' in txt and 'I CAN' in txt: t='divider'
    else: t='other'
    return (std,t)

def classify_teacher(txt):
    if 'How This Unit Supports' in txt: return ('','legend')
    std=next((s for s in STDS if s in txt),'')
    if 'KEY VOCABULARY' in txt or 'Word Wall' in txt: t='vocab'
    elif 'PRIMARY SOURCE ANALYSIS' in txt: t='source'
    elif 'STUDENT ACTIVITY' in txt: t='activity'
    elif 'CHECK FOR UNDERSTANDING' in txt: t='check'
    elif 'ANSWER' in txt and 'REVEAL' in txt.upper(): t='reveal'
    elif 'CONFIDENCE CHECK' in txt: t='confidence'
    elif 'QUICK REVIEW' in txt: t='review'
    elif 'WRAP-UP' in txt: t='wrapup'
    elif 'DIRECT INSTRUCTION' in txt: t='di'
    elif 'SOCIAL STUDIES PRACTICES' in txt or ('TN ACADEMIC STANDARD' in txt and 'I CAN' in txt): t='divider'
    elif 'Key Figure' in txt or 'PEOPLE WHO SHAPED' in txt.upper(): t='figures'
    elif 'GUIDED PRACTICE' in txt: t='guided'
    elif 'Hook' in txt or 'Think About It' in txt: t='hook'
    else: t='other'
    return (std,t)

def analyze(path,cls):
    prs=Presentation(path); rows=[]
    for i,s in enumerate(prs.slides,1):
        std,t=cls(slide_text(s)); rows.append((i,std,t))
    return rows,len(list(prs.slides))

srows,sn=analyze('working/Unit5_Student_Deck_WORKING.pptx',classify_student)
trows,tn=analyze('working/Unit5_Teacher_Deck_WORKING.pptx',classify_teacher)

def rng(nums):
    if not nums: return '—'
    nums=sorted(nums); return f"{nums[0]}" if len(nums)==1 else f"{nums[0]}–{nums[-1]}"
def nums_for(rows,std,types):
    return [i for (i,st,t) in rows if st==std and t in types]

print(f"Student deck slides: {sn} | Teacher deck slides: {tn}\n")
cw=[]
for std in STDS:
    s_all=nums_for(srows,std,{'divider','vocab','di','source','perspectives','progress'})
    t_all=nums_for(trows,std,{'divider','review','confidence','hook','di','figures','source','vocab','guided','activity','check','reveal','wrapup','other'})
    row={
     'standard':std,
     'student_range':rng(s_all),
     'teacher_range':rng(t_all),
     'student_direct_teaching':rng(nums_for(srows,std,{'di'})),
     'teacher_direct_teaching':rng(nums_for(trows,std,{'di'})),
     'student_vocabulary':rng(nums_for(srows,std,{'vocab'})),
     'teacher_vocabulary':rng(nums_for(trows,std,{'vocab'})),
     'student_primary_source':rng(nums_for(srows,std,{'source'})),
     'teacher_primary_source':rng(nums_for(trows,std,{'source'})),
     'teacher_activity_responsechoice':rng(nums_for(trows,std,{'activity'})),
     'student_progress_check':rng(nums_for(srows,std,{'progress'})),
     'teacher_check_dok2':rng(nums_for(trows,std,{'check'})),
    }
    cw.append(row)

# markdown table (Cornell references use the STUDENT / Lean deck)
hdr=['Standard','Student Deck range','Teacher Deck range','Cornell direct-teaching (Student)','Vocabulary (Student)','Primary source (Student)','Response Choice / Activity (Teacher)','Progress Check (Student)']
lines=['| '+' | '.join(hdr)+' |','|'+'|'.join(['---']*len(hdr))+'|']
for r in cw:
    lines.append('| '+' | '.join([r['standard'],r['student_range'],r['teacher_range'],r['student_direct_teaching'],
      r['student_vocabulary'],r['student_primary_source'],r['teacher_activity_responsechoice'],r['student_progress_check']])+' |')
table='\n'.join(lines)
print(table)
json.dump({'student_slides':sn,'teacher_slides':tn,'rows':cw,
          'student_classification':srows,'teacher_classification':trows},
         open('analysis/crosswalk.json','w'),indent=1)
open('analysis/Unit5_Slide_Crosswalk.md','w').write(
  '# Unit 5 — Verified Lean→Full Slide Crosswalk (from FINAL edited decks)\n\n'
  f'Student (Lean) Deck: **{sn} slides** · Teacher (Full) Deck: **{tn} slides**. '
  'Every number read from the final PPTX. Cornell references use the Student (Lean) Deck.\n\n'+table+'\n')
print('\nwrote analysis/Unit5_Slide_Crosswalk.md + crosswalk.json')
