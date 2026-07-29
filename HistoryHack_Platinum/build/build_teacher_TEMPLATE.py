#!/usr/bin/env python3
"""Unit 5 Teacher (Full) Deck — revision pass on a WORKING COPY. Originals untouched."""
import json, zipfile, shutil, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
SRC='source_decks/Unit5_Teacher_Deck.pptx'; OUT='working/Unit5_Teacher_Deck_WORKING.pptx'
NAVY='1A2332'; NAVY2='243047'; GOLD='C9A84C'; GOLDBR='F9A825'; RED='C62828'; GREEN='1B5E20'
WHITE='FFFFFF'; SLATEL='C7D4E0'; MUTE='5C6470'; CREAM='F7F4EC'; INK='1F2430'; PANEL='2C384F'; LINEC='DAD3C2'
HEAD='Georgia'; BODY='Calibri'; LABEL='Trebuchet MS'; W=13.333; H=7.5
FOOT='U.S. History Hack  ·  TroopToTeacher Technologies LLC  ·  Aligned to TN Academic Standards & TCAP EOC'
def C(h): return RGBColor.from_string(h)
STANDARDS=['US.39','US.40','US.41','US.42','US.43','US.44']
ACTIVITY=[17,34,51,68,85,102]; WRAPUP=[20,37,54,71,88,105]; CFU=[18,35,52,69,86,103]; REVEAL=[19,36,53,70,87,104]
WORDWALL=[15,32,49,66,83,100]; PRIMARY=[14,31,48,65,82,99]
PERM={'US.39':[0,2,1,3],'US.40':[1,0,2,3],'US.41':[0,2,3,1],'US.42':[0,1,2,3],'US.43':[1,0,2,3],'US.44':[0,2,1,3]}
def new_correct(std): return PERM[std].index(1)   # index in 0..3
IMGALT={  # rich alt for known slugs
 'bank-run-american-union-bank-1931':'Black-and-white photograph of a dense crowd of depositors outside the American Union Bank during a 1931 bank run.',
 'stock-market-crash-crowd-nyse-1929':'Black-and-white photograph of a large crowd on Broad Street outside the New York Stock Exchange, October 1929.',
 'dust-bowl-farmer-fence-cimarron-1936':'Black-and-white FSA photograph by Arthur Rothstein of a farmer raising a fence post from drifted Dust Bowl sand, 1936.',
 'dust-bowl-buried-machinery-dallas-sd-1936':'Black-and-white USDA photograph of farm machinery half-buried by wind-blown Dust Bowl soil, South Dakota, 1936.',
 'depression-soup-kitchen-breadline-nyc-1930s':'Black-and-white photograph of a long breadline of unemployed men in New York City, early 1930s.',
 'migrant-mother-lange-1936':"Black-and-white FSA photograph by Dorothea Lange, 'Migrant Mother,' 1936.",
 'hooverville-shantytown-1930s':'Black-and-white photograph of a Depression-era Hooverville of makeshift scrap-wood and tin shacks.',
 'herbert-hoover-portrait':'Formal studio portrait of President Herbert Hoover.',
 'hoover-dam-construction-1930s':'Black-and-white photograph of Hoover Dam under construction in the early 1930s.',
 'bonus-army-camp-anacostia-1932':'Black-and-white Harris & Ewing photograph of the Bonus Army camp at Anacostia Flats, Washington, D.C., 1932.',
 'bonus-army-eviction-1932':'Black-and-white 1932 U.S. Army Signal Corps photograph of the forced eviction of the Bonus Army camp.',
 'ccc-workers-laying-stone-1933':'Black-and-white photograph of Civilian Conservation Corps workers laying stone on a hillside, c. 1933.',
 'fdr-signing-social-security-1935':'Black-and-white photograph of President Franklin D. Roosevelt signing the Social Security Act, 1935.',
 'fdr-fireside-chat-1930s':'Black-and-white photograph of President Franklin D. Roosevelt delivering a radio fireside chat in the 1930s.',
 'tva-douglas-dam-construction-1942':'Color FSA/OWI photograph of Douglas Dam under construction in Tennessee, 1942.',
 'wpa-work-pays-america-poster-1936':"Color WPA screenprint poster, 'Work Pays America! Prosperity.'",
 'social-security-poster-1935':'Color 1935 federal poster promoting the new Social Security program.',
}
def title_words(slug):
    w=slug.replace('-',' ').replace('_',' ')
    return w[:1].upper()+w[1:]
def alt_for(slug):
    if slug in IMGALT: return IMGALT[slug]
    if 'portrait' in slug: return 'Portrait of '+title_words(slug.replace('-portrait','').replace('portrait',''))+'.'
    if 'poster' in slug: return 'Poster: '+title_words(slug.replace('-poster',''))+'.'
    if 'cartoon' in slug: return 'Political cartoon: '+title_words(slug)+'.'
    return 'Photograph: '+title_words(slug)+'.'

def rect(s,x,y,w,h,fill,line=None,lw=1.0,rounded=False):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.shadow.inherit=False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=C(fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=C(line); sh.line.width=Pt(lw)
    return sh
def txt(s,x,y,w,h,runs,size,color,font=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    if isinstance(runs,str): runs=[[(runs,{})]]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if sp: p.line_spacing=sp
        for t,opt in para:
            r=p.add_run(); r.text=t; f=r.font
            f.size=Pt(opt.get('size',size)); f.name=opt.get('font',font); f.bold=opt.get('bold',bold); f.italic=opt.get('italic',italic); f.color.rgb=C(opt.get('color',color))
    return tb
def footer(s): txt(s,0.5,7.08,11.0,0.3,FOOT,9,SLATEL,font=LABEL)
def new_slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[0])
    for ph in list(s.placeholders): ph._element.getparent().remove(ph._element)
    return s
def opt_text_shapes(s,xlo,xhi):
    xs=[sh for sh in s.shapes if sh.has_text_frame and xlo<=Emu(sh.left or 0).inches<=xhi and 2.5<Emu(sh.top or 0).inches<6.0 and len(sh.text_frame.text.strip())>3]
    return sorted(xs,key=lambda sh:Emu(sh.top).inches)
def set_text_keepfmt(shape,newtext):
    tf=shape.text_frame; p=tf.paragraphs[0]
    if p.runs: p.runs[0].text=newtext
    else: p.add_run().text=newtext
    for r in p.runs[1:]: r.text=''
def set_color(shape,hexcol,bold=None):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb=C(hexcol)
            if bold is not None: r.font.bold=bold

shutil.copy(SRC,OUT); prs=Presentation(OUT)
prs.core_properties.title='U.S. History Hack — Unit 5 Course Standard Teacher Deck'
prs.core_properties.author='TroopToTeacher Technologies LLC'
sl=list(prs.slides)

# (1) alt text + 9pt floor + autofit on Word Wall & Primary Source panels
for idx,s in enumerate(sl,1):
    for sh in s.shapes:
        if sh.shape_type==13:
            cn=sh._element.xpath('.//p:cNvPr')[0]; d=cn.get('descr') or ''
            if '/' in d or d.endswith('.jpg'):
                cn.set('descr', alt_for(d.split('/')[-1].replace('.jpg','')))
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is not None and r.font.size<Pt(9): r.font.size=Pt(9)
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() and not (Emu(sh.left or 0).inches>11.9 and Emu(sh.top or 0).inches>6.9):
            sh.text_frame.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


# (1b) image-rights: remove the Acme/AP FDR-signing photo (unverified PD basis) on orig slide 79
s79=sl[78]
for sh in list(s79.shapes):
    onright = (Emu(sh.left or 0).inches>7.4 and 1.5<Emu(sh.top or 0).inches<5.6)
    if sh.shape_type==13 and onright: sh._element.getparent().remove(sh._element)          # the picture
    elif sh.has_text_frame and ('Acme' in sh.text_frame.text or sh.text_frame.text.strip()=='ZOOM IN'): sh._element.getparent().remove(sh._element)
    elif sh.shape_type==1 and onright and (not sh.has_text_frame or not sh.text_frame.text.strip()): sh._element.getparent().remove(sh._element)  # orphan image frame
# widen the left instruction/body boxes to use the freed space
for sh in s79.shapes:
    if sh.has_text_frame and Emu(sh.left or 0).inches<1.2 and Emu(sh.width or 0).inches<8 and sh.text_frame.text.strip():
        sh.width=Inches(11.9)

# (2) label vocabulary text replacements (on-slide)
GLOBAL=[('(EL & writing support)','· LANGUAGE SUPPORT (for every learner)')]
ACT_REPL=[('DIFFERENTIATION','SUPPORT OPTION  ·  EXTENSION'),('Entry:','SUPPORT OPTION:'),('Honors:','EXTENSION:')]
WRAP_REPL=[('★ HONORS / EXTENSION','★ EXTENSION'),('HONORS / EXTENSION','EXTENSION')]
def repl_runs(s,pairs,notes=False):
    for sh in s.shapes:
        if not sh.has_text_frame: continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                for a,b in pairs:
                    if a in r.text: r.text=r.text.replace(a,b)
    if notes and s.has_notes_slide:
        tf=s.notes_slide.notes_text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                for a,b in pairs:
                    if a in r.text: r.text=r.text.replace(a,b)
for s in sl: repl_runs(s,GLOBAL)
for n in ACTIVITY: repl_runs(sl[n-1],ACT_REPL)
# (3) fixed-track language in speaker notes too
NOTE_REPL=[('HONORS / EXTENSION','EXTENSION'),(', for early finishers or advanced learners',' — optional deeper challenge, open to any student'),('for early finishers or advanced learners','optional deeper challenge, open to any student')]
for n in WRAPUP: repl_runs(sl[n-1],WRAP_REPL,notes=False); repl_runs(sl[n-1],NOTE_REPL,notes=True)

# (4) factual corrections
def replace_anytext(s,a,b,notes=True):
    for sh in s.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if a in r.text: r.text=r.text.replace(a,b)
    if notes and s.has_notes_slide:
        for p in s.notes_slide.notes_text_frame.paragraphs:
            for r in p.runs:
                if a in r.text: r.text=r.text.replace(a,b)
for s in sl:
    replace_anytext(s,'The largest, in Central Park, housed 15,000 people.','Some of the largest — such as the one in St. Louis — housed around a thousand residents.')
    replace_anytext(s,'25 cents/hour), 40-hour work week','25 cents/hour), a 44-hour work week dropping to 40 hours by 1940')
    replace_anytext(s,'maximum 40-hour work week','maximum 44-hour work week (dropping to 40 hours by 1940)')
    # rights claim on closing slide
    replace_anytext(s,'All images public domain.','Images are public-domain or openly licensed and credited to their hosting archive.')
    replace_anytext(s,'All images are public-domain, captioned, and linked','Images are public-domain or openly licensed, captioned, and linked',notes=True)

# (5) RESPONSE CHOICE strip on activity slides
for n in ACTIVITY:
    txt(sl[n-1],2.55,1.72,10.2,0.44,[[('RESPONSE CHOICE  ',{'bold':True,'color':RED,'font':LABEL}),('write it  ·  say & record it  ·  diagram & label it  ',{'bold':True,'color':INK}),('— all three demonstrate the same learning target.',{'italic':True,'color':MUTE})]],11,INK,font=BODY,anchor=MSO_ANCHOR.MIDDLE)
# (6) decision band + scaffold fading on wrap-ups
for n in WRAPUP:
    s=sl[n-1]; rect(s,0.6,6.42,12.1,0.58,'FBEFD0',line=GOLD,lw=1.0,rounded=True)
    txt(s,0.8,6.46,11.8,0.52,[[('TEACHER DECISION BAND   ',{'bold':True,'color':RED,'font':LABEL,'size':9.5}),('PROGRESS CHECK → identify barrier → select support → reteach → recheck → extend',{'bold':True,'color':INK,'size':10})],[('SCAFFOLD FADING   ',{'bold':True,'color':RED,'font':LABEL,'size':9}),('Guided → Light → independent on evidence; reintroduce for a new barrier. Required IEP/504 accommodations do not fade unless the authorized team changes them.',{'color':INK,'size':9})]],9.5,INK,font=BODY,anchor=MSO_ANCHOR.MIDDLE,sp=1.0)

# (7) de-bias: reorder CFU + reveal options in sync with student deck (move correct answer position)
for std,cfu_n,rev_n in zip(STANDARDS,CFU,REVEAL):
    perm=PERM[std]; nc=new_correct(std); letter='ABCD'[nc]
    # CFU: reorder the 4 option texts
    cs=opt_text_shapes(sl[cfu_n-1],1.2,1.5)
    if len(cs)==4:
        old=[sh.text_frame.text for sh in cs]
        for pos in range(4): set_text_keepfmt(cs[pos], old[perm[pos]])
    # REVEAL: reorder texts, move highlight+check+letter
    rs=sl[rev_n-1]
    ts=opt_text_shapes(rs,1.2,1.5)
    if len(ts)==4:
        old=[sh.text_frame.text for sh in ts]
        for pos in range(4): set_text_keepfmt(ts[pos], old[perm[pos]])
        for pos in range(4): set_color(ts[pos], NAVY if pos==nc else WHITE, bold=(pos==nc))
        rows=sorted([sh for sh in rs.shapes if sh.shape_type==1 and abs(Emu(sh.width or 0).inches-12.1)<0.2 and 2.5<Emu(sh.top or 0).inches<6.0],key=lambda sh:Emu(sh.top).inches)
        for pos,rw in enumerate(rows[:4]):
            rw.fill.solid(); rw.fill.fore_color.rgb=C(GOLD if pos==nc else PANEL)
        # move check glyph to correct row y
        chk=[sh for sh in rs.shapes if sh.has_text_frame and sh.text_frame.text.strip()=='✓' and Emu(sh.left or 0).inches>11.5]
        if chk and len(rows)>=4: chk[0].top=rows[nc].top
        # update "Correct answer: X"
        replace_anytext(rs,'Correct answer: B','Correct answer: '+letter,notes=True)
        replace_anytext(rs,'Correct: B','Correct: '+letter,notes=True)
    # notes on CFU
    replace_anytext(sl[cfu_n-1],'Correct answer: B','Correct answer: '+letter,notes=True)

# (8) legend slide -> position 3
def build_legend(prs):
    s=new_slide(prs); rect(s,0,0,W,H,NAVY); rect(s,0,0,W,0.16,GOLD)
    txt(s,0.7,0.55,12,0.4,'UNIVERSAL DESIGN FOR LEARNING · MTSS  —  IMPLEMENTATION LEGEND',12,GOLDBR,font=LABEL,bold=True)
    txt(s,0.7,0.95,12,0.8,'How This Unit Supports Every Learner',32,WHITE,font=HEAD,bold=True)
    labels=[('CORE PATH','The grade-level lesson every student does — same TN standard, same rigor, same ceiling.'),('SUPPORT OPTION','An optional scaffold that lowers a barrier, not the goal. Available to anyone.'),('LANGUAGE SUPPORT','Vocabulary, pronunciation, Spanish, and sentence frames for multilingual learners — open to all.'),('RESPONSE CHOICE','Students show understanding their way: write it · say & record it · diagram & label it. Same target.'),('PROGRESS CHECK','A quick DOK-2 check and confidence rating that routes the next instructional step.'),('EXTENSION','An optional deeper challenge (higher DOK). The standard stays the same.')]
    x0,y0,cw,ch,gx,gy=0.7,2.0,5.95,1.15,0.25,0.18
    for i,(lab,desc) in enumerate(labels):
        x=x0+(i%2)*(cw+gx); y=y0+(i//2)*(ch+gy)
        rect(s,x,y,cw,ch,NAVY2,line=GOLD,lw=0.75,rounded=True)
        txt(s,x+0.25,y+0.12,cw-0.5,0.3,lab,13,GOLDBR,font=LABEL,bold=True)
        txt(s,x+0.25,y+0.46,cw-0.5,ch-0.55,desc,11.5,SLATEL,font=BODY,sp=1.03)
    rect(s,0.7,5.95,12.0,0.95,'11202F',line=GOLDBR,lw=1.0,rounded=True)
    txt(s,0.95,6.06,11.5,0.78,[[('Options are available by design and support access without lowering the learning goal.  ',{'bold':True,'color':WHITE}),('Supports work alongside, never in place of, required IEP or 504 accommodations.',{'color':GOLDBR,'bold':True})]],13,WHITE,font=BODY,anchor=MSO_ANCHOR.MIDDLE,sp=1.05)
    footer(s); return s

def build_rep(prs):
    s=new_slide(prs); rect(s,0,0,W,H,CREAM); rect(s,0,0,W,0.16,RED)
    txt(s,0.6,0.42,12,0.3,'US.41 & US.43 · WHOSE STORY THE NEW DEAL REACHED',12,RED,font=LABEL,bold=True)
    txt(s,0.6,0.78,12.1,0.7,'Reached — and Left Out',30,NAVY,font=HEAD,bold=True)
    blocks=[('Black Americans','The AAA\u2019s crop-reduction payments went to landowners, pushing many Black sharecroppers and tenant farmers off the land, and Social Security (1935) at first excluded farm and domestic workers \u2014 jobs held by most Black workers. Yet FDR\u2019s informal \u201cBlack Cabinet,\u201d including Mary McLeod Bethune and Robert Weaver, advised the administration, and Black voters shifted toward the Democratic Party.'),
      ('Mexican Americans','Amid Depression job competition, \u201cMexican Repatriation\u201d (about 1929\u20131936) pressured an estimated 400,000 to more than one million people of Mexican descent \u2014 many of them U.S.-born citizens \u2014 to leave for Mexico. Mexican and Filipino farmworkers in California often faced lower pay and worse conditions than white \u201cOkie\u201d migrants.'),
      ('Indigenous Peoples','The Indian Reorganization Act of 1934 (the \u201cIndian New Deal\u201d), led by John Collier, ended the Dawes-era allotment policy, restored some tribal self-government and land, and supported Native culture \u2014 a major reversal, though its results were mixed and debated. A CCC Indian Division employed thousands.')]
    cw=3.95; gap=0.2; x0=0.6; y=1.95; ch=4.2; hh=0.5
    for i,(h,b) in enumerate(blocks):
        x=x0+i*(cw+gap)
        rect(s,x,y,cw,ch,WHITE,line=LINEC,lw=1.0); rect(s,x,y,cw,hh,NAVY)
        txt(s,x+0.18,y,cw-0.36,hh,h,14,WHITE,font=BODY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+0.18,y+hh+0.14,cw-0.36,ch-hh-0.28,b,11.5,INK,font=BODY,sp=1.06)
    txt(s,0.6,6.5,12.1,0.4,[[('History Hack-authored instructional synthesis grounded in the standard historical record \u2014 not a primary source. ',{'italic':True,'color':MUTE}),('Verified public-domain images and primary sources for these histories are a pending content addition.',{'italic':True,'color':RED})]],9.5,MUTE,font=BODY,sp=1.0)
    if s.has_notes_slide or True:
        s.notes_slide.notes_text_frame.text='REPRESENTATION (US.41 & US.43). Use to complicate the New Deal narrative: who was included, excluded, and organized. Pair with the Three Perspectives routine. Facts are standard historical record; sourced images/primary sources are a pending content commission.'
    footer(s); return s
legend=build_legend(prs)
lst=prs.slides._sldIdLst; kids=list(lst); lst.remove(kids[-1]); lst.insert(2,kids[-1])
rep=build_rep(prs)                       # appended at very end
kids=list(lst); rep_el=kids[-1]; lst.remove(rep_el); lst.insert(len(list(lst)), rep_el)
kids=list(lst); rep_el=kids[-1]; lst.remove(rep_el); lst.insert(len(list(lst))-1, rep_el)  # before unit-close

# (9) page numbers -> file position; add to legend
for i,s in enumerate(prs.slides,1):
    num=None
    for sh in s.shapes:
        if sh.has_text_frame and Emu(sh.left or 0).inches>11.9 and Emu(sh.top or 0).inches>6.9 and sh.text_frame.text.strip().isdigit(): num=sh; break
    if num is not None:
        r=num.text_frame.paragraphs[0].runs
        (r[0] if r else num.text_frame.paragraphs[0].add_run()).text=str(i)
    else:
        txt(s,12.23,7.08,0.6,0.3,str(i),9,SLATEL,font=LABEL,align=PP_ALIGN.RIGHT)

prs.save(OUT)
z=zipfile.ZipFile(OUT); files={it.filename:z.read(it.filename) for it in z.infolist()}; infos=z.infolist(); z.close()
ct=files['[Content_Types].xml'].decode()
if 'Extension="jpg"' not in ct:
    files['[Content_Types].xml']=ct.replace('<Default Extension="jpeg" ContentType="image/jpeg"/>','<Default Extension="jpeg" ContentType="image/jpeg"/><Default Extension="jpg" ContentType="image/jpeg"/>').encode()
if 'docProps/app.xml' in files:
    files['docProps/app.xml']=re.sub(r'<Slides>\d+</Slides>','<Slides>108</Slides>',files['docProps/app.xml'].decode()).encode()
with zipfile.ZipFile(OUT,'w',zipfile.ZIP_DEFLATED) as o:
    for it in infos: o.writestr(it, files[it.filename])
print('SAVED',OUT,'slides:',len(list(Presentation(OUT).slides)),'| CFU/reveal correct letters:',{s:'ABCD'[new_correct(s)] for s in STANDARDS})
