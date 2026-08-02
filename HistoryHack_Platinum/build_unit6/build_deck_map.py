#!/usr/bin/env python3
"""build_deck_map.py — emit the reproducible activity<->slide MAP for a unit.

This is the deck-side half of the print-bundle standard (skill:
history-hack-workbook-print-bundle). The workbook build reads the map to print
EXACT deck slide numbers on each activity header, and the deck stamp step reads
it to put "In your workbook . Activity N" back-refs on the matching slides. Both
sides come from THIS one map, so the numbers can never drift.

The workbook (student-facing) references the STUDENT deck slides — the deck a
student can actually navigate for review. Teacher-deck slides are recorded too.

Usage: python3 build_deck_map.py TEACHER_DECK.pptx STUDENT_DECK.pptx OUT_MAP.json
"""
import json
import re
import sys
from pptx import Presentation

# role keyword -> canonical role
ROLE_KEYS = [
    ("KEY VOCABULARY", "KEY VOCABULARY"),
    ("DIRECT INSTRUCTION", "DIRECT INSTRUCTION"),
    ("PRIMARY SOURCE", "SOURCE IT FIRST"),
    ("SOURCE IT FIRST", "SOURCE IT FIRST"),
    ("THREE PERSPECTIVES", "THREE PERSPECTIVES"),
    ("STUDENT ACTIVITY", "THREE PERSPECTIVES"),
    ("CHECK FOR UNDERSTANDING", "PROGRESS CHECK"),
    ("PROGRESS CHECK", "PROGRESS CHECK"),
]
# canonical role -> workbook activity numbers it backs
ROLE_TO_ACTIVITY = {
    "KEY VOCABULARY": [1, 2],
    "DIRECT INSTRUCTION": [3, 4],
    "SOURCE IT FIRST": [5],
    "PROGRESS CHECK": [6],
    "THREE PERSPECTIVES": [7],
}
STD_RE = re.compile(r"(US\.\d+)")


def first_title(sl):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().split("\n")[0]
    return ""


def role_of(title):
    up = title.upper()
    for kw, canon in ROLE_KEYS:
        if kw in up:
            return canon
    return None


def extract(path):
    """Return {US.xx: {role: [slide_numbers]}} for a deck, handling both the
    teacher convention (each slide titled 'US.xx . ROLE') and the student
    convention (a bare 'US.xx' divider starts a block, role slides follow)."""
    prs = Presentation(path)
    mp = {}
    cur = None
    for i, sl in enumerate(prs.slides, 1):
        title = first_title(sl)
        head = title.strip().split("\n")[0]
        m = STD_RE.match(head)
        bare_divider = bool(re.match(r"^US\.\d+\s*$", head))
        if m:
            cur = m.group(1)
            mp.setdefault(cur, {})
            if bare_divider:
                continue  # student-deck divider: no role on this slide
        r = role_of(title)
        if r and cur:
            mp[cur].setdefault(r, []).append(i)
    return mp


def activity_slides(student_map):
    """Per standard: {activityN: [student-deck slide numbers]} for the workbook."""
    out = {}
    for std, roles in student_map.items():
        acts = {}
        for role, slides in roles.items():
            for act in ROLE_TO_ACTIVITY.get(role, []):
                acts.setdefault(f"Activity {act}", [])
                acts[f"Activity {act}"].extend(slides)
        for k in acts:
            acts[k] = sorted(set(acts[k]))
        out[std] = acts
    return out


def main():
    teacher, student, outp = sys.argv[1], sys.argv[2], sys.argv[3]
    tmap = extract(teacher)
    smap = extract(student)
    result = {
        "reference_deck": "student",
        "teacher": tmap,
        "student": smap,
        "activity_slides": activity_slides(smap),
        "role_to_activity": {k: v for k, v in ROLE_TO_ACTIVITY.items()},
    }
    json.dump(result, open(outp, "w"), ensure_ascii=False, indent=1)
    print(f"wrote {outp}: teacher={len(tmap)} std, student={len(smap)} std")
    # DI-parity report (teacher DI count vs student DI count)
    print("DI parity (teacher/student):")
    for std in sorted(tmap):
        td = len(tmap.get(std, {}).get("DIRECT INSTRUCTION", []))
        sd = len(smap.get(std, {}).get("DIRECT INSTRUCTION", []))
        flag = "" if td == sd else "  <-- MISMATCH"
        print(f"  {std}: teacher {td} / student {sd}{flag}")


if __name__ == "__main__":
    main()
