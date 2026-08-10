#!/usr/bin/env python3
"""
Apply accessibility QC fixes to the Unit 1 student deck (python-pptx):
 F1  recolor #C7D4E0 subtitle text -> Heritage Blue #1F3A5F ONLY where it sits on
     a light background (keeps the one instance on the dark #232F40 panel).
 F2  descriptive alt text on every image, built from the slide's caption /
     provenance / topic (replaces "ZOOM IN" / "image.jpg").
 F3  "WE DO" chip text white -> navy (matches the other chips).
 F4  tag Spanish "ES:" runs lang="es".
 F6  set presentation Title metadata.
Usage: fix_student_deck_a11y.py IN.pptx OUT.pptx
"""
import sys, re
from pptx import Presentation
from pptx.dml.color import RGBColor

NAVY = RGBColor(0x1F,0x3A,0x5F)
A="http://schemas.openxmlformats.org/drawingml/2006/main"
P="http://schemas.openxmlformats.org/presentationml/2006/main"

def lum(h):
    r,g,b=[int(h[i:i+2],16)/255 for i in (0,2,4)]
    f=lambda c: c/12.92 if c<=0.03928 else ((c+0.055)/1.055)**2.4
    return 0.2126*f(r)+0.7152*f(g)+0.0722*f(b)
def contrast(a,b):
    L1,L2=lum(a),lum(b); hi,lo=max(L1,L2),min(L1,L2); return (hi+.05)/(lo+.05)

def run_hex(run):
    try:
        if run.font.color and run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
    except Exception: pass
    return None

def shape_fill_hex(sh):
    try:
        if sh.fill.type == 1:  # MSO_FILL.SOLID
            return str(sh.fill.fore_color.rgb).upper()
    except Exception: pass
    return None

YEAR=re.compile(r"\b1[89]\d\d\b")
INST=("Library of Congress","National Archives","HathiTrust","Public Domain",
      "Tennessee State Library","U.S. Reports","Statutes","Census","Freeman",
      "Scribner","Macmillan","North American Review","Smithsonian","NARA")
SKIP=("U.S. History Hack","ZOOM IN","RESPONSE CHOICE","Say it, see it","FOLLOW ALONG",
      "Then you try","Watch how","Now you complete","commit before","MODELED","TEACHER THINK",
      "CAUSE →","Trace ONE","Trace one","complete two more","historian returns","end of each reading")

def looks_caption(s):
    s=s.strip()
    if not s or any(k in s for k in SKIP): return False
    if s.isdigit(): return False
    # a real caption carries provenance (a year or a hosting institution),
    # or is a descriptive phrase about a map/photo/image
    if bool(YEAR.search(s)) or any(k in s for k in INST): return True
    return bool(re.search(r"\b(map|photograph|photo|portrait|cartoon|version of)\b", s, re.I))

def build_alt(slide):
    # collect (top, left, text, is_topic)
    lines=[]
    topic=None
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t=sh.text_frame.text.strip()
        if not t: continue
        if re.match(r"US\.0\d\s*·", t) and topic is None:
            topic=re.sub(r"^US\.0\d\s*·\s*","",t).strip()
        lines.append((sh.top or 0, sh.left or 0, t))
    alts=[]
    pics=[sh for sh in slide.shapes if sh.shape_type==13]  # PICTURE
    for pic in pics:
        py=(pic.top or 0)+(pic.height or 0)
        cands=sorted([(abs(ty-py), ty, tx, tt) for (ty,tx,tt) in lines
                      if looks_caption(tt) and ty>=(pic.top or 0)], key=lambda z:z[0])
        if cands:
            alt="Primary source: "+cands[0][3]
        elif topic:
            alt=f"Primary-source image — {topic}."
        else:
            alt="Historical primary-source image, Unit 1 (U.S. History)."
        alts.append((pic, alt[:250]))
    return alts

def main(src,dst):
    prs=Presentation(src)
    f1=f2=f3=f4=0
    for slide in prs.slides:
        infos=[(sh, sh.left or 0, sh.top or 0, sh.width or 0, sh.height or 0, shape_fill_hex(sh))
               for sh in slide.shapes]
        def bg_behind(i, cx, cy):
            bg="FFFFFF"
            for j in range(i):
                _,l,t,w,h,rgb=infos[j]
                if rgb and l<=cx<=l+w and t<=cy<=t+h: bg=rgb
            return bg
        # alt text
        for pic,alt in build_alt(slide):
            cNvPr=pic._element._nvXxPr.cNvPr
            cNvPr.set("descr", alt); f2+=1
        # text fixes
        for i,(sh,l,t,w,h,rgb) in enumerate(infos):
            if not sh.has_text_frame: continue
            cx,cy=l+w/2,t+h/2
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    hexc=run_hex(run); txt=run.text.strip()
                    if hexc=="C7D4E0":
                        if contrast("C7D4E0", bg_behind(i,cx,cy))<4.5:
                            run.font.color.rgb=NAVY; f1+=1
                    if txt=="WE DO":
                        run.font.color.rgb=NAVY; f3+=1
                    if txt.startswith("ES:") or txt.startswith("ES ·") or re.search(r"[áéíóúñ]", txt):
                        run._r.get_or_add_rPr().set("lang","es"); f4+=1
    prs.core_properties.title="U.S. History Hack — Unit 1 Student Deck (America 250)"
    prs.save(dst)
    print(f"F1 C7D4E0->navy (light bg): {f1} | F2 alt set: {f2} | F3 WE DO->navy: {f3} | F4 lang=es: {f4}")

if __name__=="__main__":
    main(sys.argv[1], sys.argv[2])
